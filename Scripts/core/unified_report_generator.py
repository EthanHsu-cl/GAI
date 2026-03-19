import json, yaml, logging, sys, re, tempfile, os
from datetime import datetime
from pathlib import Path
from dataclasses import dataclass, field
from typing import Dict, List, Optional, Union
from concurrent.futures import ThreadPoolExecutor
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm, Inches, Pt
from pptx.enum.text import PP_ALIGN

from config_loader import get_app_base_path, get_resource_path, get_core_path

try:
    import cv2
except ImportError:
    cv2 = None

try:
    from tqdm import tqdm
    HAS_TQDM = True
except ImportError:
    HAS_TQDM = False
    # Fallback progress tracker
    class tqdm:
        def __init__(self, iterable=None, total=None, desc=None, **kwargs):
            self.iterable = iterable
            self.total = total or (len(iterable) if iterable else 0)
            self.desc = desc
            self.n = 0
        
        def __iter__(self):
            for item in self.iterable:
                yield item
                self.update()
        
        def update(self, n=1):
            self.n += n
            if self.total and self.n % max(1, self.total // 10) == 0:
                logger.info(f"{self.desc}: {self.n}/{self.total} ({100*self.n//self.total}%)")
        
        def close(self):
            pass

logging.basicConfig(level=logging.INFO, format='%(message)s')
logger = logging.getLogger(__name__)

@dataclass
class MediaPair:
    """Universal media pair for all API types"""
    source_file: str
    source_path: Path
    api_type: str
    
    # Generated content
    generated_paths: List[Path] = field(default_factory=list)
    reference_paths: List[Path] = field(default_factory=list)
    
    # API-specific fields
    effect_name: str = ""
    category: str = ""
    prompt: str = ""
    
    # Video-specific (for Runway)
    source_video_path: Optional[Path] = None
    
    # Multi-image support (for Nano Banana dual-image mode)
    additional_source_paths: List[Path] = field(default_factory=list)
    
    # Metadata
    metadata: Dict = field(default_factory=dict)
    ref_metadata: Dict = field(default_factory=dict)

    # Status
    failed: bool = False
    ref_failed: bool = False
    
    @property
    def primary_generated(self) -> Optional[Path]:
        """Get the primary generated file"""
        return self.generated_paths[0] if self.generated_paths else None
    
    @property
    def primary_reference(self) -> Optional[Path]:
        """Get the primary reference file"""
        return self.reference_paths[0] if self.reference_paths else None


class UnifiedReportGenerator:
    """
    Unified Report Generator - Single implementation for all APIs
    No old methods included - completely unified approach
    """
    
    # ================== LAYOUT CONSTANTS ==================
    # Common layout configurations for easier maintenance
    
    # 2-Media Layout: Two large squares side-by-side (edge-to-edge)
    LAYOUT_2_MEDIA = {
        'positions': [(0.42, 2.15, 16, 16), (17.44, 2.15, 16, 16)],
        'metadata_position': (35, 0, 7.29, 3.06),
        'metadata_reference_position': (2.32, 15.26, 7.29, 3.06),
    }
    
    # 3-Media Layout: Three medium squares (for comparison/multi-image)
    LAYOUT_3_MEDIA = {
        'positions': [(2.59, 3.26, 10, 10), (13, 3.26, 10, 10), (23.41, 3.26, 10, 10)],
        'metadata_position': (2.32, 15.24, 7.29, 3.06),
    }

    # 3-Media Stacked Layout: Two sources stacked left, big generated right
    # Heights are dynamically adjusted per-slide based on actual aspect ratios.
    # These are the default fallback positions when aspect ratios are unavailable.
    # Metadata uses same top-right position as LAYOUT_2_MEDIA for consistency.
    LAYOUT_3_MEDIA_STACKED = {
        'positions': [(0.42, 2.15, 16, 7), (0.42, 9.55, 16, 7), (17.44, 2.15, 16, 16)],
        'metadata_position': (35, 0, 7.29, 3.06),
        'metadata_reference_position': (35, 0, 7.29, 3.06),
        'override_positions': True,
        'media_labels': ['Source Image', 'Source Video', None],
    }
    
    # File extension constants
    IMAGE_EXTS = {'.jpg', '.jpeg', '.png', '.webp'}
    VIDEO_EXTS = {'.mp4', '.mov', '.avi'}
    SUPPORTED_IMG_FORMATS = {'.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff', '.tif'}
    METADATA_EXTS = {'.json'}
    
    def __init__(self, api_name: str, config_file: str = None):
        self.api_name = api_name
        # Auto-detect YAML or JSON config files
        if not config_file:
            yaml_path = Path(f"config/batch_{api_name}_config.yaml")
            json_path = Path(f"config/batch_{api_name}_config.json")
            if yaml_path.exists():
                self.config_file = str(yaml_path)
            elif json_path.exists():
                self.config_file = str(json_path)
            else:
                self.config_file = f"config/batch_{api_name}_config.yaml"  # Default to YAML
        else:
            self.config_file = config_file
        self.config = {}
        self.report_definitions = {}
        
        # Caches for performance
        self._ar_cache = {}
        self._frame_cache = {}
        self._tempfiles_to_cleanup = []  # Track temporary files from format conversions
        self._normalize_cache = {}  # Cache for normalize_key operations
        self._extract_key_cache = {}  # Cache for video key extraction
        
        # Smart batching configuration
        self._batch_size = 50  # Process 50 items at a time
        self._max_workers = 4  # Default thread pool size
        self._show_progress = HAS_TQDM  # Use tqdm if available
        
        # API display names mapping (used across multiple methods)
        # Note: For kling, the actual display name is set dynamically based on config model
        self._api_display_names = {
            'kling': 'Kling',  # Will be updated based on model in config
            'kling_effects': 'Kling Effects',
            'kling_endframe': 'Kling Endframe',  # Will be updated based on model in config
            'kling_ttv': 'Kling TTV',  # Will be updated based on model in config
            'nano_banana': 'Nano Banana',
            'runway': 'Runway',
            'vidu_effects': 'Vidu Effects',
            'vidu_reference': 'Vidu Reference',
            'genvideo': 'GenVideo',
            'pixverse': 'Pixverse',
            'wan': 'Wan 2.2',
            'dreamactor': 'DreamActor',
            'kling_motion': 'Kling Motion',
            'veo': 'Veo',
            'veo_itv': 'Veo ITV'
        }
        
        # Load configurations
        self.load_config()
        self.load_report_definitions()
        
        # Update Kling display name based on model in config
        if self.api_name in ['kling', 'kling_endframe', 'kling_ttv', 'kling_motion']:
            self._update_kling_display_name()
    
    # ================== UNIFIED CONFIGURATION SYSTEM ==================
    
    def get_slide_config(self):
        """Get API-specific slide configuration"""
        # All APIs use section dividers and group by effect_name by default
        base_config = {
            'use_section_dividers': True,
            'group_by': 'effect_name',
        }
        configs = {
            'runway': {
                **base_config,
                'media_types': ['source', 'source_video', 'generated'],
                **self.LAYOUT_3_MEDIA_STACKED,
                'title_format': 'Generation {index}: {source_file}',
                'metadata_fields': ['prompt', 'reference_image', 'source_video', 'model', 'processing_time_seconds', 'success'],
                'error_handling': 'video_fallback'
            },
            'nano_banana': {
                **base_config,
                'media_types': ['source', 'generated'],  # Default 2-media layout
                **self.LAYOUT_2_MEDIA,
                'title_format': '❌ GENERATION FAILED',
                'title_show_only_if_failed': True,
                'metadata_fields': ['response_id', 'additional_images_used', 'success', 'attempts', 'processing_time_seconds'],
                'use_comparison': False,
                'supports_multi_image': True,
                # Alternative 3-media layout (used when additional_images are present)
                'media_types_3': ['source', 'additional_source', 'generated'],
                'positions_3': self.LAYOUT_3_MEDIA['positions']
            },
            'vidu_effects': {
                **base_config,
                'media_types': ['source', 'generated'],
                **self.LAYOUT_2_MEDIA,
                'title_format': 'Generation {index}: {source_file}',
                'metadata_fields': ['effect_name', 'category', 'task_id', 'processing_time_seconds', 'duration', 'success'],
            },
            'vidu_reference': {
                **base_config,
                'media_types': ['source', 'generated'],
                **self.LAYOUT_2_MEDIA,
                'title_format': 'Generation {index}: {source_file}',
                'metadata_fields': ['effect_name', 'reference_count', 'task_id', 'processing_time_seconds', 'duration', 'success'],
                'supports_multi_image': True,
            },
            'genvideo': {
                **base_config,
                'media_types': ['source', 'generated', 'reference'],
                **self.LAYOUT_2_MEDIA,
                'title_format': 'GenVideo {index}: {source_file}',
                'metadata_fields': ['model', 'quality', 'processing_time_seconds', 'success', 'img_prompt'],
            },
            'pixverse': {
                **base_config,
                'media_types': ['source', 'generated'],
                **self.LAYOUT_2_MEDIA,
                'title_format': 'pixverse_{index}_{source_file}',
                'metadata_fields': ['effect_name', 'video_id', 'processing_time_seconds', 'success'],
            },
            'kling': {
                **base_config,
                'media_types': ['source', 'generated', 'reference'],
                **self.LAYOUT_2_MEDIA,
                'title_format': 'Generation {index}: {source_file}',
                'metadata_fields': ['task_id', 'model', 'prompt', 'processing_time_seconds', 'success'],
            },
            'kling_effects': {
                **base_config,
                'media_types': ['source', 'generated'],
                **self.LAYOUT_2_MEDIA,
                'title_format': 'Generation {index}: {source_file}',
                'metadata_fields': ['effect_name', 'custom_effect', 'video_id', 'task_id', 'duration', 'processing_time_seconds', 'success'],
            },
            'kling_endframe': {
                **base_config,
                'media_types': ['source', 'source_video', 'generated'],
                **self.LAYOUT_3_MEDIA_STACKED,
                'title_format': 'Generation {index}: {source_file}',
                'metadata_fields': ['start_image', 'end_image', 'generation_number', 'task_id', 'model', 'prompt', 'processing_time_seconds', 'success'],
                'error_handling': 'video_fallback',
                'media_labels': ['Start Frame', 'End Frame', None],
            },
            'wan': {
                **base_config,
                'media_types': ['source', 'source_video', 'generated'],
                **self.LAYOUT_3_MEDIA_STACKED,
                'title_format': 'Generation {index}: {source_file}',
                'metadata_fields': ['source_image', 'source_video', 'animation_mode', 'prompt', 'processing_time_seconds', 'success'],
                'error_handling': 'video_fallback'
            },
            'dreamactor': {
                **base_config,
                'media_types': ['source', 'source_video', 'generated'],
                **self.LAYOUT_3_MEDIA_STACKED,
                'title_format': 'Generation {index}: {source_file}',
                'metadata_fields': ['source_image', 'source_video', 'task_id', 'time_taken', 'status_code', 'processing_time_seconds', 'success'],
                'error_handling': 'video_fallback'
            },
            'kling_motion': {
                **base_config,
                'media_types': ['source', 'source_video', 'generated'],
                **self.LAYOUT_3_MEDIA_STACKED,
                'title_format': 'Generation {index}: {source_file}',
                'metadata_fields': ['source_image', 'source_video', 'model', 'character_orientation', 'mode', 'video_id', 'task_id', 'processing_time_seconds', 'success'],
                'error_handling': 'video_fallback'
            },
            'veo': {
                **base_config,
                'media_types': ['generated'],
                'positions': [(17.44, 2.15, 16, 16)],  # Single centered video
                **self.LAYOUT_2_MEDIA,  # Use standard 2-media metadata position (top-right)
                'title_format': 'Generation {index}: {style_name}',
                'metadata_fields': ['style_name', 'generation_number', 'model_id', 'duration_seconds', 'aspect_ratio', 'resolution', 'processing_time_seconds', 'success'],
                'use_section_dividers': False,
                'group_by': None
            },
            'veo_itv': {
                **base_config,
                'media_types': ['source', 'generated'],
                **self.LAYOUT_2_MEDIA,
                'title_format': 'Generation {index}: {source_file}',
                'metadata_fields': ['source_image', 'generation_number', 'model_id', 'duration_seconds', 'aspect_ratio', 'resolution', 'processing_time_seconds', 'success'],
            },
            'kling_ttv': {
                **base_config,
                'media_types': ['prompt', 'generated'],  # Prompt text box + video
                'positions': [(0.42, 2.15, 16, 16), (17.44, 2.15, 16, 16)],  # Prompt left, video right
                **self.LAYOUT_2_MEDIA,  # Use standard 2-media metadata position (top-right)
                'title_format': 'Generation {index}: {style_name}',
                'metadata_fields': ['style_name', 'generation_number', 'model', 'mode', 'duration', 'ratio', 'cfg', 'processing_time_seconds', 'success'],
                'use_section_dividers': False,
                'group_by': None
            }
        }
        return configs.get(self.api_name, configs['kling'])
    
    # ================== UNIFIED SLIDE CREATION ENGINE ==================
    
    def create_slides(self, ppt, pairs, template_loaded, use_comparison=False):
        """Universal slide creation for all APIs"""
        slide_config = self.get_slide_config()
        grouped_pairs = self.group_pairs_if_needed(pairs, slide_config)
        
        slide_index = 1
        for group_name, group_pairs in grouped_pairs.items():
            # Add section divider if needed
            if slide_config.get('use_section_dividers') and group_name != 'default':
                self.create_section_divider_slide(ppt, group_name, template_loaded)
            
            # Create individual slides
            for pair in group_pairs:
                self.create_universal_slide(ppt, pair, slide_index, template_loaded, 
                                          use_comparison, slide_config)
                slide_index += 1
    
    def create_universal_slide(self, ppt, pair, index, template_loaded,
                              use_comparison, slide_config):
        """Create a single slide for any API using configuration"""
        # Adjust media types for nano_banana based on whether multi-image mode is active
        if self.api_name == 'nano_banana':
            if pair.additional_source_paths:
                # Check if any additional source is a video
                has_video_source = any(
                    p.suffix.lower() in self.VIDEO_EXTS 
                    for p in pair.additional_source_paths if p
                )
                
                if has_video_source:
                    # 1 image + 1 video: use 3-media layout (source, video, generated)
                    slide_config = slide_config.copy()
                    slide_config['media_types'] = slide_config.get('media_types_3', ['source', 'additional_source', 'generated'])
                    slide_config['positions'] = slide_config.get('positions_3', self.LAYOUT_3_MEDIA['positions'])
                # else: All images (2+ images) - use default 2-media layout with composite
                # The composite will be created in get_media_path_and_type()
        
        # Create slide
        if template_loaded and len(ppt.slides) >= 4:
            slide = ppt.slides.add_slide(ppt.slides[3].slide_layout)
            self.handle_template_slide(slide, pair, index, use_comparison, slide_config)
        else:
            slide = ppt.slides.add_slide(ppt.slide_layouts[6])
            self.handle_manual_slide(slide, pair, index, use_comparison, slide_config)
    
    def _format_title(self, pair, index, title_format, show_only_if_failed):
        """Format title for slide based on configuration"""
        if show_only_if_failed and not pair.failed:
            return None
        
        # Build format kwargs with all possible placeholders
        format_kwargs = {
            'index': index,
            'source_file': pair.source_file or '',
            'failure_status': "❌ GENERATION FAILED" if pair.failed else "",
            'style_name': pair.effect_name or '',  # effect_name is used for style_name
        }
        
        return title_format.format(**format_kwargs)
    
    def _compute_stacked_positions(self, pair, base_positions):
        """Dynamically adjust stacked layout heights based on source aspect ratios.

        Redistributes the vertical space between the two stacked source boxes
        so each source fills its allocated area as much as possible, regardless
        of whether sources are landscape, portrait, or square.

        Args:
            pair: MediaPair with source_path and source_video_path.
            base_positions: List of 3 position tuples from LAYOUT_3_MEDIA_STACKED.

        Returns:
            List of 3 adjusted position tuples (x, y, w, h).
        """
        if len(base_positions) != 3:
            return base_positions

        source_path = pair.source_path
        source_video_path = pair.source_video_path

        # Get actual aspect ratios (width / height). Default landscape if unreadable.
        ar1 = (self.get_aspect_ratio(source_path, False)
               if source_path and source_path.exists() else 16 / 9)
        is_second_video = (source_video_path.suffix.lower() in self.VIDEO_EXTS
                           if source_video_path and source_video_path.exists() else True)
        ar2 = (self.get_aspect_ratio(source_video_path, is_second_video)
               if source_video_path and source_video_path.exists() else 16 / 9)

        # Layout parameters derived from the base positions
        x = base_positions[0][0]         # Left x (0.42)
        y_start = base_positions[0][1]   # Top y (2.15)
        w = base_positions[0][2]         # Box width (16)
        total_h = base_positions[2][3]   # Match generated height (16)
        gap = 0.4                        # Vertical gap between stacked items
        min_h = 3.0                      # Minimum height for any source box

        usable_h = total_h - gap

        # For a box of width w, the height needed to display at full width is w/ar.
        # Allocate proportionally so each source gets the height it "wants".
        h1_ideal = w / ar1
        h2_ideal = w / ar2
        total_ideal = h1_ideal + h2_ideal

        h1 = usable_h * (h1_ideal / total_ideal)
        h2 = usable_h * (h2_ideal / total_ideal)

        # Enforce minimum heights
        if h1 < min_h:
            h1 = min_h
            h2 = usable_h - h1
        elif h2 < min_h:
            h2 = min_h
            h1 = usable_h - h2

        y2 = y_start + h1 + gap

        return [
            (x, y_start, w, round(h1, 2)),
            (x, round(y2, 2), w, round(h2, 2)),
            base_positions[2],  # Generated box stays unchanged
        ]

    def handle_template_slide(self, slide, pair, index, use_comparison, slide_config):
        """Handle slide creation with template placeholders (optimized)"""
        # Update title placeholder
        title = self._format_title(
            pair, index, 
            slide_config.get('title_format', 'Generation {index}: {source_file}'),
            slide_config.get('title_show_only_if_failed', False)
        )
        
        if title:
            # Optimized: Find title placeholder directly
            title_ph = next((p for p in slide.placeholders if p.placeholder_format.type == 1), None)
            if title_ph:
                title_ph.text = title
                if pair.failed and title_ph.text_frame.paragraphs:
                    title_ph.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
        
        if slide_config.get('override_positions', False):
            # Stacked layout: remove template placeholders, use custom positions
            for ph in list(slide.placeholders):
                if ph.placeholder_format.type in {6, 7, 8, 13, 18, 19}:
                    try:
                        ph._element.getparent().remove(ph._element)
                    except Exception:
                        pass
            
            positions = self._compute_stacked_positions(
                pair, slide_config.get('positions', [])
            )
            media_types = slide_config.get('media_types', ['source', 'generated'])
            media_labels = slide_config.get('media_labels', [])
            
            for idx, (pos, media_type) in enumerate(zip(positions, media_types)):
                # Add label above media item if specified
                if idx < len(media_labels) and media_labels[idx]:
                    label_y = max(pos[1] - 0.45, 0)
                    label_box = slide.shapes.add_textbox(
                        Cm(pos[0]), Cm(label_y), Cm(6), Cm(0.45)
                    )
                    label_box.text_frame.text = media_labels[idx]
                    label_box.text_frame.paragraphs[0].font.size = Pt(8)
                    label_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
                    label_box.text_frame.paragraphs[0].font.bold = True
                
                media_path, is_video = self.get_media_path_and_type(pair, media_type)
                self.add_media_universal(slide, pos, media_path, is_video, slide_config, pair, media_type)
        else:
            # Standard: use template placeholder positions
            phs = sorted([p for p in slide.placeholders 
                         if p.placeholder_format.type in {6, 7, 8, 13, 18, 19}],
                        key=lambda x: getattr(x, 'left', 0))
            
            media_types = slide_config.get('media_types', ['source', 'generated'])
            
            for i, (ph, media_type) in enumerate(zip(phs, media_types)):
                media_path, is_video = self.get_media_path_and_type(pair, media_type)
                self.add_media_universal(slide, ph, media_path, is_video, slide_config, pair, media_type)
        
        # Add metadata
        self.add_metadata_universal(slide, pair, slide_config, use_comparison)
    
    def handle_manual_slide(self, slide, pair, index, use_comparison, slide_config):
        """Handle slide creation without template"""
        # Add title if needed
        title = self._format_title(
            pair, index,
            slide_config.get('title_format', 'Generation {index}: {source_file}'),
            slide_config.get('title_show_only_if_failed', False)
        )
        
        if title:
            tb = slide.shapes.add_textbox(Cm(2), Cm(1), Cm(20), Cm(2))
            tb.text_frame.text = title
            tb.text_frame.paragraphs[0].font.size = Pt(20)
            if pair.failed:
                tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
        
        # Add media using positions
        positions = slide_config.get('positions', [(2.59, 3.26, 12.5, 12.5), (18.78, 3.26, 12.5, 12.5)])
        if slide_config.get('override_positions', False):
            positions = self._compute_stacked_positions(pair, positions)
        media_types = slide_config.get('media_types', ['source', 'generated'])
        media_labels = slide_config.get('media_labels', [])
        
        for idx, (pos, media_type) in enumerate(zip(positions, media_types)):
            # Add label above media item if specified
            if idx < len(media_labels) and media_labels[idx]:
                label_y = max(pos[1] - 0.45, 0)
                label_box = slide.shapes.add_textbox(
                    Cm(pos[0]), Cm(label_y), Cm(6), Cm(0.45)
                )
                label_box.text_frame.text = media_labels[idx]
                label_box.text_frame.paragraphs[0].font.size = Pt(8)
                label_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
                label_box.text_frame.paragraphs[0].font.bold = True
            
            media_path, is_video = self.get_media_path_and_type(pair, media_type)
            self.add_media_universal(slide, pos, media_path, is_video, slide_config, pair, media_type)
        
        # Add metadata
        self.add_metadata_universal(slide, pair, slide_config, use_comparison)
    
    def get_media_path_and_type(self, pair, media_type):
        """Get media path and determine if it's video.
        
        For nano_banana with multiple source images (either random_source_selection or
        standard multi-image mode), creates a composite grid image for the 'source' media type
        when all additional sources are images.
        """
        # For multi-image APIs: create composite from all source images
        # nano_banana: source + additional images from Additional folder
        # vidu_reference: source + reference images from Reference folder
        if (media_type == 'source' and 
            self.api_name in ('nano_banana', 'vidu_reference') and 
            pair.additional_source_paths):
            # Check if all additional sources are images (not videos)
            all_images = all(
                p.suffix.lower() not in self.VIDEO_EXTS 
                for p in pair.additional_source_paths if p
            )
            if all_images:
                # Determine which images to include in composite
                if pair.metadata.get('random_source_selection'):
                    # random_source_selection: additional_source_paths already contains all images
                    all_sources = list(pair.additional_source_paths)
                else:
                    # Standard multi-image: combine source_path with additional_source_paths
                    all_sources = [pair.source_path] + list(pair.additional_source_paths)
                
                composite_path = self._create_source_composite(all_sources)
                if composite_path:
                    return Path(composite_path), False
        
        path_map = {
            'source': pair.source_path,
            'source_video': pair.source_video_path,
            'additional_source': pair.additional_source_paths[0] if pair.additional_source_paths else None,
            'generated': pair.primary_generated,
            'reference': pair.primary_reference,
            'prompt': None  # Special type for text-to-video prompt display
        }
        
        path = path_map.get(media_type)
        # For text-to-video APIs (veo, kling_ttv), generated content is always video
        is_video = (media_type in ['source_video', 'generated'] and self.api_name in ['veo', 'kling_ttv']) or \
                   (media_type == 'source_video') or \
                   (path and path.suffix.lower() in self.VIDEO_EXTS)
        
        return path, is_video
    
    def group_pairs_if_needed(self, pairs, slide_config):
        """Group pairs if needed for section dividers and apply sorting"""
        group_by = slide_config.get('group_by')
        if group_by:
            grouped = {}
            for pair in pairs:
                key = getattr(pair, group_by, 'default')
                grouped.setdefault(key, []).append(pair)
            
            # Sort each group
            for key in grouped:
                grouped[key] = self._sort_pairs(grouped[key])
            return grouped
        else:
            # Sort the entire list if not grouped
            return {'default': self._sort_pairs(pairs)}
    
    def _sort_pairs(self, pairs: List[MediaPair]) -> List[MediaPair]:
        """Sort pairs based on API type
        
        Failed slides are always placed first within each section.
        Then:
        - For combination APIs (Wan, Runway): Group by video/reference, then sort by source filename
        - For nano_banana with random_source_selection: Sort by number of source images used (min to max)
        - For other APIs: Sort by source filename only
        """
        if not pairs:
            return pairs
        
        # Check if this is nano_banana iteration mode (random_source_selection)
        if self.api_name == 'nano_banana':
            # Check if any pair has random_source_selection metadata
            has_random_selection = any(
                p.metadata.get('random_source_selection') for p in pairs
            )
            if has_random_selection:
                # Sort by failed status first (failed slides at start),
                # then by number of source images used (min to max),
                # then by iteration index for stable ordering within same count
                def get_sort_key(pair):
                    # Primary: failed status (not failed so True=failed comes first)
                    failed_priority = not pair.failed
                    
                    # Secondary: number of all_images_used
                    all_images = pair.metadata.get('all_images_used', [])
                    num_images = len(all_images) if all_images else len(pair.additional_source_paths)
                    
                    # Tertiary: iteration index for stable ordering
                    iteration_idx = pair.metadata.get('_iteration_index', 0)
                    
                    return (failed_priority, num_images, iteration_idx)
                
                return sorted(pairs, key=get_sort_key)
        
        if self.api_name in ['wan', 'runway', 'dreamactor', 'kling_motion']:
            # Combination APIs: Failed slides first, then group by source video,
            # then sort within groups by source file
            def get_sort_key(pair):
                # Primary key: failed status (not failed so True=failed comes first)
                failed_priority = not pair.failed
                
                # Secondary key: source video name (or source file if no video)
                if pair.source_video_path:
                    video_name = pair.source_video_path.name
                else:
                    video_name = ""
                
                # Tertiary key: source file name
                source_name = pair.source_file or ""
                
                return (failed_priority, video_name, source_name)
            
            return sorted(pairs, key=get_sort_key)
        else:
            # Other APIs: Failed slides first, then sort by source filename
            return sorted(pairs, key=lambda p: (not p.failed, p.source_file or ""))
    
    # ================== UNIFIED MEDIA SYSTEM ==================
    
    def ensure_supported_img_format(self, img_path):
        """Convert any unsupported image format to PNG for PowerPoint compatibility.
        
        Also applies EXIF orientation to fix rotated images from cameras/phones.
        """
        p = Path(img_path)
        
        # PowerPoint natively supports: jpg, jpeg, png, bmp, gif, tiff, tif
        # We'll convert everything else (webp, svg, mpo, etc.) to PNG
        needs_conversion = p.suffix.lower() not in self.SUPPORTED_IMG_FORMATS
        
        # Always try to open and check the actual format, not just extension
        # Some files like .jpg might actually be MPO format
        try:
            with Image.open(p) as im:
                # Check if EXIF orientation needs to be applied
                has_exif_orientation = False
                try:
                    exif = im.getexif()
                    # Orientation tag is 274
                    if exif and exif.get(274, 1) != 1:
                        has_exif_orientation = True
                except Exception:
                    pass
                
                # Check actual image format - MPO and other exotic formats need conversion
                actual_format = im.format
                if actual_format in ('MPO', 'WEBP', 'SVG', 'HEIC', 'HEIF', 'AVIF'):
                    needs_conversion = True
                    logger.info(f"Detected {actual_format} format in {p.name}, will convert")
                
                # Also convert if EXIF orientation needs to be applied
                if has_exif_orientation:
                    needs_conversion = True
                
                if needs_conversion:
                    # Apply EXIF orientation to fix rotated images
                    try:
                        im = ImageOps.exif_transpose(im)
                    except Exception:
                        pass  # If EXIF transpose fails, continue with original
                    
                    # Convert to RGB mode (removes alpha for formats that don't support it well)
                    # Use RGBA for formats that might have transparency
                    mode = 'RGBA' if im.mode in ('RGBA', 'LA', 'P') else 'RGB'
                    rgb_im = im.convert(mode)
                    tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
                    rgb_im.save(tmp.name, 'PNG')
                    tmp.close()
                    self._tempfiles_to_cleanup.append(tmp.name)
                    logger.info(f"Converted {actual_format or p.suffix} to PNG: {p.name}")
                    return tmp.name
                else:
                    # Format is supported and no orientation fix needed, return as-is
                    return str(img_path)
                    
        except Exception as e:
            logger.error(f"Failed to process image {p.name}: {e}")
            # If we can't even open it, try converting anyway as a fallback
            if needs_conversion:
                logger.error(f"Unable to convert {p.name}, returning original path")
            return str(img_path)
        
        return str(img_path)
    
    def _create_source_composite(self, image_paths: List[Path], cell_size: int = 400) -> Optional[str]:
        """Create a composite grid image from multiple source images.
        
        Creates an N×N grid where each cell maintains the source image's aspect ratio
        with letterboxing/pillarboxing using a light gray background.
        
        Args:
            image_paths: List of paths to source images.
            cell_size: Size of each grid cell in pixels (cells are square).
        
        Returns:
            str: Path to temporary composite image file, or None on failure.
        """
        if not image_paths:
            return None
        
        if len(image_paths) == 1:
            # Single image - no composite needed
            return str(image_paths[0])
        
        # Determine grid dimensions
        count = len(image_paths)
        if count <= 2:
            cols, rows = 2, 1
        elif count <= 4:
            cols, rows = 2, 2
        elif count <= 6:
            cols, rows = 3, 2
        elif count <= 9:
            cols, rows = 3, 3
        elif count <= 12:
            cols, rows = 4, 3
        else:
            cols, rows = 4, 4  # Max 16 images
        
        # Light gray background color
        bg_color = (240, 240, 240)
        
        # Create composite canvas
        composite_width = cols * cell_size
        composite_height = rows * cell_size
        composite = Image.new('RGB', (composite_width, composite_height), bg_color)
        
        try:
            for idx, img_path in enumerate(image_paths[:cols * rows]):
                col = idx % cols
                row = idx // cols
                
                # Calculate cell position
                cell_x = col * cell_size
                cell_y = row * cell_size
                
                try:
                    with Image.open(img_path) as img:
                        # Apply EXIF orientation to fix rotated images
                        try:
                            img = ImageOps.exif_transpose(img)
                        except Exception:
                            pass  # If EXIF transpose fails, continue with original
                        
                        # Convert to RGB if needed
                        if img.mode in ('RGBA', 'LA', 'P'):
                            # Create background for transparent images
                            bg = Image.new('RGB', img.size, bg_color)
                            if img.mode == 'P':
                                img = img.convert('RGBA')
                            bg.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                            img = bg
                        elif img.mode != 'RGB':
                            img = img.convert('RGB')
                        
                        # Calculate scaled size maintaining aspect ratio
                        img_width, img_height = img.size
                        img_ratio = img_width / img_height
                        
                        if img_ratio > 1:
                            # Landscape: fit width, letterbox height
                            new_width = cell_size
                            new_height = int(cell_size / img_ratio)
                        else:
                            # Portrait or square: fit height, pillarbox width
                            new_height = cell_size
                            new_width = int(cell_size * img_ratio)
                        
                        # Resize image
                        img_resized = img.resize((new_width, new_height), Image.LANCZOS)
                        
                        # Calculate position to center in cell
                        paste_x = cell_x + (cell_size - new_width) // 2
                        paste_y = cell_y + (cell_size - new_height) // 2
                        
                        composite.paste(img_resized, (paste_x, paste_y))
                        
                except Exception as e:
                    logger.warning(f"Failed to add image to composite: {img_path.name}: {e}")
                    # Draw placeholder for failed image
                    continue
            
            # Save composite to temp file
            tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
            composite.save(tmp.name, 'PNG', quality=95)
            tmp.close()
            self._tempfiles_to_cleanup.append(tmp.name)
            
            logger.debug(f"Created composite image: {cols}x{rows} grid with {len(image_paths)} images")
            return tmp.name
            
        except Exception as e:
            logger.error(f"Failed to create composite image: {e}")
            return None
    
    def _convert_unsupported_formats_batch(self, image_paths):
        """Convert multiple unsupported image formats in parallel for major performance gain"""
        if not image_paths:
            return {}
        
        def convert_one(path):
            try:
                converted = self.ensure_supported_img_format(path)
                return (path, converted)
            except Exception as e:
                logger.warning(f"Failed to convert {path}: {e}")
                return (path, str(path))
        
        with ThreadPoolExecutor(max_workers=4) as executor:
            results = executor.map(convert_one, image_paths)
        
        return dict(results)
    
    def _load_json_batch(self, json_files_dict):
        """Load multiple JSON files in parallel - 40-50% faster metadata loading"""
        if not json_files_dict:
            return {}
        
        def load_one(key_path_tuple):
            key, path = key_path_tuple
            try:
                with open(path, 'r', encoding='utf-8') as f:
                    return (key, json.load(f))
            except Exception as e:
                logger.warning(f"Failed to load JSON {path.name}: {e}")
                return (key, {})
        
        items = list(json_files_dict.items())
        with ThreadPoolExecutor(max_workers=8) as executor:
            if self._show_progress and len(items) > 20:
                pbar = tqdm(total=len(items), desc="Loading metadata", unit="files")
                results = []
                for result in executor.map(load_one, items):
                    results.append(result)
                    pbar.update()
                pbar.close()
            else:
                results = list(executor.map(load_one, items))
        
        return dict(results)
    
    def add_media_universal(self, slide, placeholder_or_pos, media_path, is_video, slide_config, pair=None, media_type=None):
        """Universal media addition for all APIs with webp conversion
        
        Note: Always removes placeholders and adds media as new shapes for consistency.
        This ensures proper behavior when mixing images and videos on the same slide.
        """
        # Handle both placeholder objects and manual positions
        if hasattr(placeholder_or_pos, 'left'):
            # It's a placeholder or shape - extract position and remove it
            l, t, w, h = (placeholder_or_pos.left, placeholder_or_pos.top,
                         placeholder_or_pos.width, placeholder_or_pos.height)
            try:
                placeholder_or_pos._element.getparent().remove(placeholder_or_pos._element)
            except:
                pass  # Placeholder might already be removed
        else:
            # It's a position tuple
            if len(placeholder_or_pos) == 4:
                l, t, w, h = [Cm(x) for x in placeholder_or_pos]
            else:
                l, t, w, h = Cm(placeholder_or_pos[0]), Cm(placeholder_or_pos[1]), Cm(10), Cm(10)
        
        # Special handling for 'prompt' media type - display as text box
        if media_type == 'prompt' and pair and pair.metadata:
            prompt_text = pair.metadata.get('prompt', 'No prompt available')
            box = slide.shapes.add_textbox(l, t, w, h)
            box.text_frame.text = prompt_text
            box.text_frame.word_wrap = True
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(245, 245, 245)
            box.line.color.rgb = RGBColor(200, 200, 200)
            box.line.width = Pt(1)
            for para in box.text_frame.paragraphs:
                para.font.size = Pt(11)
                para.alignment = PP_ALIGN.LEFT
            return
        
        if media_path and Path(media_path).exists():
            try:
                # Calculate aspect ratio and positioning for all media types
                ar = self.get_aspect_ratio(Path(media_path), is_video)
                sw, sh = (w, w/ar) if ar > w/h else (h*ar, h)
                fl, ft = l + (w - sw)/2, t + (h - sh)/2
                
                if is_video:
                    # Extract first frame for video poster
                    first_frame_path = self.extract_first_frame(Path(media_path))
                    if first_frame_path and Path(first_frame_path).exists():
                        slide.shapes.add_movie(str(media_path), fl, ft, sw, sh,
                                             poster_frame_image=first_frame_path)
                    else:
                        slide.shapes.add_movie(str(media_path), fl, ft, sw, sh)
                else:
                    # Convert any unsupported image format to PNG if needed
                    converted_path = self.ensure_supported_img_format(media_path)
                    slide.shapes.add_picture(str(converted_path), fl, ft, sw, sh)
            except Exception as e:
                self.add_error_box(slide, l, t, w, h, f"Failed to load media: {e}", pair)
        else:
            # Extract failure message from metadata if available
            error_msg = self.get_failure_message(pair) if pair else None
            self.add_error_box(slide, l, t, w, h, error_msg or "Media not found", pair)
    
    def get_failure_message(self, pair):
        """Extract failure message from metadata"""
        if not pair or not pair.metadata:
            return None
        
        # Check error field first (from handler)
        error_msg = pair.metadata.get('error', '')
        if error_msg:
            return f"{error_msg}"
        
        # Try to get text_responses for detailed error message
        text_responses = pair.metadata.get('text_responses', [])
        if text_responses and isinstance(text_responses, list):
            for response in text_responses:
                if isinstance(response, dict) and 'content' in response:
                    content = response['content']
                    if content and content.strip():
                        return f"Media not found\n\nAPI Response:\n{content}"
        
        # Fallback to generic message
        return "Media not found"
    
    def add_error_box(self, slide, left, top, width, height, message: str, pair=None):
        """Add error box with proper styling"""
        box = slide.shapes.add_textbox(left, top, width, height)
        box.text_frame.text = f"❌ GENERATION FAILED\n\n{message}"
        box.text_frame.word_wrap = True
        
        for para in box.text_frame.paragraphs:
            para.font.size = Pt(12)  # Slightly smaller to fit more text
            para.alignment = PP_ALIGN.CENTER
            para.font.color.rgb = RGBColor(255, 0, 0)
        
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 240, 240)
        box.line.color.rgb = RGBColor(255, 0, 0)
        box.line.width = Pt(0.5)
    
    # ================== UNIFIED METADATA SYSTEM ==================
    
    def _add_metadata_field(self, field, pair, meta_lines):
        """Add a single metadata field to meta_lines"""
        md = pair.metadata or {}
        
        field_handlers = {
            'success': lambda: f"Status: {'✓' if md.get('success', False) else '❌'}",
            'processing_time_seconds': lambda: f"Time: {md.get(field, 'N/A')}s",
            'response_id': lambda: f"Response ID: {md.get(field, 'N/A')}",
            'attempts': lambda: f"Attempts: {md.get(field, 'N/A')}",
            'task_id': lambda: f"Task ID: {md.get(field, 'N/A')}",
            'effect_name': lambda: f"Effect: {pair.effect_name}",
            'category': lambda: f"Category: {pair.category}",
            'start_image': lambda: f"Start: {md.get(field, 'N/A')}",
            'end_image': lambda: f"End: {md.get(field, 'N/A')}",
            'source_image': lambda: f"Image: {md.get(field, 'N/A')}",
            'source_video': lambda: f"Video: {md.get(field, 'N/A')}",
            'animation_mode': lambda: f"Mode: {md.get(field, 'N/A')}",
            'style_name': lambda: f"Style: {md.get(field, 'N/A')}",
            'model_id': lambda: f"Model: {md.get(field, 'N/A')}",
            'duration_seconds': lambda: f"Duration: {md.get(field, 'N/A')}s",
            'aspect_ratio': lambda: f"Aspect: {md.get(field, 'N/A')}",
            'resolution': lambda: f"Resolution: {md.get(field, 'N/A')}",
            'model': lambda: f"Model: {md.get(field, 'N/A')}",
            'mode': lambda: f"Mode: {md.get(field, 'N/A')}",
            'duration': lambda: f"Duration: {md.get(field, 'N/A')}s",
            'ratio': lambda: f"Ratio: {md.get(field, 'N/A')}",
            'cfg': lambda: f"CFG: {md.get(field, 'N/A')}",
        }
        
        # Special handlers
        if field == 'additional_images_used':
            add_imgs = md.get(field, [])
            if add_imgs:
                text = add_imgs[0] if len(add_imgs) == 1 else ', '.join(add_imgs)
                meta_lines.append(f"Additional: {text}")
        elif field == 'generation_number':
            gen_num = md.get('generation_number')
            total_gens = md.get('total_generations', 1)
            if gen_num and total_gens > 1:
                meta_lines.append(f"Generation: {gen_num}/{total_gens}")
        elif field in ['prompt', 'img_prompt']:
            value = md.get(field, 'N/A')
            text = f"{str(value)[:60]}..." if len(str(value)) > 60 else str(value)
            meta_lines.append(f"Prompt: {text}")
        elif field in field_handlers:
            meta_lines.append(field_handlers[field]())
        else:
            # Generic field
            value = md.get(field, 'N/A')
            display_name = field.replace('_', ' ').title()
            meta_lines.append(f"{display_name}: {value}")
    
    def add_metadata_universal(self, slide, pair, slide_config, use_comparison=False):
        """Universal metadata addition for all APIs"""
        metadata_fields = slide_config.get('metadata_fields', [])
        metadata_pos = slide_config.get('metadata_reference_position', (2.32, 15.26, 7.29, 3.06)) if use_comparison else slide_config.get('metadata_position', (5.19, 15.99, 7.29, 3.06))
        
        meta_lines = []
        
        # Build metadata lines based on configuration
        for field in metadata_fields:
            self._add_metadata_field(field, pair, meta_lines)
        
        # Add source file name for some APIs
        if self.api_name in ['nano_banana', 'genvideo', 'pixverse', 'kling_endframe']:
            meta_lines.insert(0, f"File: {pair.source_file}")
        
        if not meta_lines:
            meta_lines = ["No metadata available"]
        
        # Add metadata box
        box = slide.shapes.add_textbox(Cm(metadata_pos[0]), Cm(metadata_pos[1]),
                                       Cm(metadata_pos[2]), Cm(metadata_pos[3]))
        box.text_frame.text = "\n".join(meta_lines)
        box.text_frame.word_wrap = True
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        for para in box.text_frame.paragraphs:
            para.font.size = Pt(10)
    
    def create_section_divider_slide(self, ppt, effect_name, template_loaded):
        """Create section divider slide for effects"""
        if template_loaded and len(ppt.slides) >= 2:
            slide = ppt.slides.add_slide(ppt.slides[1].slide_layout)
            # Set title in placeholder
            for p in slide.placeholders:
                if p.placeholder_format.type == 1:  # Title placeholder
                    p.text = f"{effect_name}"
                    if p.text_frame.paragraphs:
                        p.text_frame.paragraphs[0].font.size = Pt(48)
                        p.text_frame.paragraphs[0].font.bold = True
                        p.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    break
        else:
            slide = ppt.slides.add_slide(ppt.slide_layouts[6])
            # Add title text box
            tb = slide.shapes.add_textbox(Cm(5), Cm(8), Cm(24), Cm(4))
            tb.text_frame.text = f"{effect_name}"
            for p in tb.text_frame.paragraphs:
                p.font.size = Pt(48)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
    
    # ================== FILE PROCESSING METHODS ==================
    
    def process_batch(self, task: Dict) -> List[MediaPair]:
        """Universal batch processing for all API types"""
        if self.api_name in ["vidu_effects", "vidu_reference", "pixverse", "kling_effects", "veo_itv"]:
            return self.process_base_folder_structure(task)
        elif self.api_name == "genvideo":
            return self.process_genvideo_batch(task)
        elif self.api_name in ["veo", "kling_ttv"]:
            return self.process_text_to_video_batch(task)
        else:
            return self.process_task_folder_structure(task)
    
    def process_task_folder_structure(self, task: Dict) -> List[MediaPair]:
        """Process APIs with individual task folders"""
        folder = Path(task['folder'])
        ref_folder = Path(task.get('reference_folder', '')) if task.get('reference_folder') else None
        slide_config = self.get_slide_config()
        if slide_config.get('override_positions', False):
            use_comparison = False
        else:
            use_comparison = task.get('use_comparison_template', False)
        
        if self.api_name == 'runway':
            return self.create_runway_media_pairs(folder, ref_folder, task, use_comparison)
        elif self.api_name in ('wan', 'dreamactor', 'kling_motion'):
            return self.create_wan_media_pairs(folder, ref_folder, task, use_comparison)
        else:
            return self.create_standard_media_pairs(folder, ref_folder, task, use_comparison)
    
    def normalize_key(self, name: str) -> str:
        """Universal key normalization - keeps aspect ratio and numbers intact (memoized)"""
        if name in self._normalize_cache:
            return self._normalize_cache[name]
        
        key = name.lower()
        # Don't remove numbers or underscores in aspect ratios like 9_16, 1_1, 16_9
        # Just remove spaces and convert to consistent format
        key = key.replace(' ', '_')
        # Keep dashes and underscores, just clean up
        result = key.strip('_')
        
        self._normalize_cache[name] = result
        return result
    
    def create_standard_media_pairs(self, folder: Path, ref_folder: Optional[Path],
                                   task: Dict, use_comparison: bool) -> List[MediaPair]:
        """Create standard media pairs for Kling/Nano Banana"""
        pairs = []
        
        # Define folder structure based on API
        if self.api_name == 'nano_banana':
            folders = {
                'source': folder / 'Source',
                'generated': folder / 'Generated_Output',
                'metadata': folder / 'Metadata'
            }
            file_pattern = 'image'
        else:  # kling or kling_endframe
            folders = {
                'source': folder / 'Source',
                'generated': folder / 'Generated_Video',
                'metadata': folder / 'Metadata'
            }
            file_pattern = 'generated'
        
        if not folders['source'].exists():
            return pairs
        
        # Get metadata with single scan and batch load first (needed to detect mode)
        _, _, metadata_files = self._scan_directory_once(folders['metadata'])
        metadata_cache = self._load_json_batch(metadata_files) if metadata_files else {}
        
        # Check if this is nano_banana iteration mode (random_source_selection)
        # In this mode, we iterate over metadata/generated files, not source files
        is_iteration_mode = (self.api_name == 'nano_banana' and 
                            any(md.get('random_source_selection') for md in metadata_cache.values()))
        
        if is_iteration_mode:
            return self._create_nano_iteration_pairs(folder, folders, metadata_cache, task)
        
        # OPTIMIZED: Single-pass directory scanning (standard mode)
        src_imgs, _, _ = self._scan_directory_once(folders['source'])
        
        # For kling_endframe, filter to only A images (start frames)
        # B images are end frames and are referenced in metadata, not source for pairs
        if self.api_name == 'kling_endframe':
            src = {k: v for k, v in src_imgs.items() if '_a_' in k or k.endswith('_a')}
        else:
            src = src_imgs
        
        # Get generated files with single scan
        out = {}
        if folders['generated'].exists():
            if self.api_name == 'nano_banana':
                gen_imgs, _, _ = self._scan_directory_once(folders['generated'])
                for key, f in gen_imgs.items():
                    if file_pattern in f.name:
                        # Split on 'image' and remove trailing underscore
                        basename = f.name.split(file_pattern)[0].rstrip('_')
                        out.setdefault(self.normalize_key(basename), []).append(f)
            else:  # kling or kling_endframe
                _, gen_vids, _ = self._scan_directory_once(folders['generated'])
                for key, f in gen_vids.items():
                    if file_pattern in f.name:
                        # Extract basename by splitting on '_generated' pattern
                        # E.g., "Name_A_generated_1.mp4" -> "Name_A"
                        basename = f.stem.split('_' + file_pattern)[0]
                        out.setdefault(self.normalize_key(basename), []).append(f)
        
        # Get metadata with single scan and batch load
        _, _, metadata_files = self._scan_directory_once(folders['metadata'])
        metadata_cache = self._load_json_batch(metadata_files) if metadata_files else {}
        
        # Get reference files
        ref_files = {}
        if use_comparison and ref_folder:
            ref_generated_folder = ref_folder / ('Generated_Output' if self.api_name == 'nano_banana' else 'Generated_Video')
            if ref_generated_folder.exists():
                for f in ref_generated_folder.iterdir():
                    if self.api_name == 'nano_banana':
                        if f.suffix.lower() in {'.jpg', '.jpeg', '.png', '.webp'} and 'image' in f.name:
                            # Split on 'image' and remove trailing underscore
                            basename = f.name.split('image')[0].rstrip('_')
                            ref_files.setdefault(self.normalize_key(basename), []).append(f)
                    else:  # kling or kling_endframe
                        if f.suffix.lower() in {'.mp4', '.mov', '.avi'} and 'generated' in f.name:
                            # Extract basename by splitting on '_generated' pattern
                            basename = f.stem.split('_generated')[0]
                            ref_files[self.normalize_key(basename)] = f
        
        # Pre-compute aspect ratios for all source images
        all_media = list(src.values()) + [p for paths in out.values() for p in paths if isinstance(paths, list)]
        if all_media:
            are_videos = {p: p.suffix.lower() in self.VIDEO_EXTS for p in all_media}
            self._compute_aspect_ratios_batch(all_media, are_videos=are_videos)
        
        # Create pairs
        for b in sorted(src.keys()):
            gen_paths = out.get(b, [])
            ref_paths = ref_files.get(b, []) if use_comparison else []

            if not isinstance(gen_paths, list):
                gen_paths = [gen_paths] if gen_paths else []
            if not isinstance(ref_paths, list):
                ref_paths = [ref_paths] if ref_paths else []

            # Determine effect_name: try metadata, then config, then folder name
            # Try to get from any available metadata first
            effect_name = None
            if gen_paths:
                # Try to get metadata from any generated file
                for gen_path in gen_paths:
                    gen_key = self.normalize_key(gen_path.stem)
                    temp_md = metadata_cache.get(gen_key, {})
                    if temp_md.get('effect_name'):
                        effect_name = temp_md['effect_name']
                        break
            
            if not effect_name:
                effect_name = self.config.get('effect') or self.config.get('effect_name')
            if not effect_name:
                # Remove leading date (e.g., '1001 ') from folder name
                m = re.match(r'^(\d{4})\s*(.+)', folder.name)
                effect_name = m.group(2) if m else folder.name

            # For Nano Banana multi-image mode: resolve additional source images from metadata
            additional_source_paths = []
            
            # Check if we have multiple generations (kling_endframe)
            # Multiple generations have files like: basename_generated_1.mp4, basename_generated_2.mp4
            if self.api_name == 'kling_endframe' and len(gen_paths) > 1:
                # Sort generated files by generation number
                sorted_gen_paths = sorted(gen_paths, key=lambda p: p.name)
                
                # Create separate MediaPair for each generation
                for gen_path in sorted_gen_paths:
                    # Get metadata for this specific generation
                    gen_key = self.normalize_key(gen_path.stem)
                    md = metadata_cache.get(gen_key, {})
                    
                    # Find end image (B frame) from metadata
                    end_image_path = None
                    end_image_name = md.get('end_image', '')
                    if end_image_name:
                        end_image_path = next(
                            (p for p in src_imgs.values() if p.name == end_image_name),
                            None
                        )
                    
                    pair = MediaPair(
                        source_file=src[b].name,
                        source_path=src[b],
                        api_type=self.api_name,
                        generated_paths=[gen_path],  # Single generated file per pair
                        reference_paths=ref_paths,
                        source_video_path=end_image_path,
                        metadata=md,
                        failed=not md.get('success', False),
                        ref_failed=use_comparison and not ref_paths,
                        effect_name=effect_name
                    )
                    pairs.append(pair)
            else:
                # Standard single generation or nano_banana
                md = metadata_cache.get(b, {})
                
                # For Nano Banana: resolve source images from metadata
                if self.api_name == 'nano_banana':
                    source_folder = folder / 'Source'
                    
                    # Check for random_source_selection mode (all_images_used in metadata)
                    if md.get('all_images_used'):
                        for img_name in md.get('all_images_used', []):
                            img_path = source_folder / img_name
                            if img_path.exists():
                                additional_source_paths.append(img_path)
                    # Fallback to additional_images_used (standard multi-image mode)
                    elif md.get('additional_images_used'):
                        additional_folder = folder / 'Additional'
                        for img_name in md.get('additional_images_used', []):
                            img_path = additional_folder / img_name
                            if img_path.exists():
                                additional_source_paths.append(img_path)

                # For kling_endframe: find end image (B frame) from metadata
                end_image_path = None
                if self.api_name == 'kling_endframe':
                    end_image_name = md.get('end_image', '')
                    if end_image_name:
                        end_image_path = next(
                            (p for p in src_imgs.values() if p.name == end_image_name),
                            None
                        )

                pair = MediaPair(
                    source_file=src[b].name,
                    source_path=src[b],
                    api_type=self.api_name,
                    generated_paths=gen_paths,
                    reference_paths=ref_paths,
                    additional_source_paths=additional_source_paths,
                    source_video_path=end_image_path,
                    metadata=md,
                    failed=not gen_paths or not md.get('success', False),
                    ref_failed=use_comparison and not ref_paths,
                    effect_name=effect_name
                )
                pairs.append(pair)
        
        return pairs
    
    def _create_nano_iteration_pairs(self, folder: Path, folders: Dict, 
                                     metadata_cache: Dict, task: Dict) -> List[MediaPair]:
        """Create media pairs for nano_banana iteration mode (random_source_selection).
        
        In iteration mode, we iterate over metadata files (which have iteration-based names
        like iter000_sourcename) and use the all_images_used field to find source images.
        
        Args:
            folder: Task folder path.
            folders: Dict with 'source', 'generated', 'metadata' paths.
            metadata_cache: Pre-loaded metadata dict.
            task: Task configuration.
        
        Returns:
            List of MediaPair objects.
        """
        pairs = []
        source_folder = folders['source']
        generated_folder = folders['generated']
        
        # Get all generated images
        gen_imgs, _, _ = self._scan_directory_once(generated_folder)
        
        # Build a map of iteration base_name -> generated files
        # e.g., "iter000_0V8A5820" -> [iter000_0V8A5820_image_5.png]
        gen_by_iteration = {}
        for key, f in gen_imgs.items():
            if 'image' in f.name:
                # Extract iteration base name: "iter000_0V8A5820_image_5.png" -> "iter000_0V8A5820"
                basename = f.name.split('image')[0].rstrip('_')
                gen_by_iteration.setdefault(self.normalize_key(basename), []).append(f)
        
        # Determine effect_name from folder
        effect_name = self.config.get('effect') or self.config.get('effect_name')
        if not effect_name:
            m = re.match(r'^(\d{4})\s*(.+)', folder.name)
            effect_name = m.group(2) if m else folder.name
        
        # Iterate over metadata entries (each represents one iteration)
        for md_key, md in sorted(metadata_cache.items()):
            if not md.get('random_source_selection'):
                continue
            
            # Get the iteration base_name from metadata
            iteration_base = md.get('_base_name', '')
            if not iteration_base:
                # Try to extract from md_key (metadata file key)
                iteration_base = md_key
            
            normalized_key = self.normalize_key(iteration_base)
            gen_paths = gen_by_iteration.get(normalized_key, [])
            
            # Get all source images used for this iteration
            additional_source_paths = []
            all_images_used = md.get('all_images_used', [])
            for img_name in all_images_used:
                img_path = source_folder / img_name
                if img_path.exists():
                    additional_source_paths.append(img_path)
            
            # Use the first source image as the "primary" source for the pair
            primary_source = additional_source_paths[0] if additional_source_paths else None
            if not primary_source:
                # Try to find source from source_image field
                source_image = md.get('source_image', '')
                if source_image:
                    primary_source = source_folder / source_image
                    if not primary_source.exists():
                        primary_source = None
            
            if not primary_source:
                logger.warning(f"No source images found for iteration {iteration_base}")
                continue
            
            pair = MediaPair(
                source_file=primary_source.name,
                source_path=primary_source,
                api_type=self.api_name,
                generated_paths=gen_paths,
                reference_paths=[],
                additional_source_paths=additional_source_paths,
                metadata=md,
                failed=not gen_paths or not md.get('success', False),
                ref_failed=False,
                effect_name=effect_name
            )
            pairs.append(pair)
        
        return pairs
    
    def create_runway_media_pairs(self, folder: Path, ref_folder: Optional[Path],
                                 task: Dict, use_comparison: bool) -> List[MediaPair]:
        """Create Runway media pairs"""
        folders = {
            'reference': folder / 'Reference',
            'source': folder / 'Source',
            'generated': folder / 'Generated_Video',
            'metadata': folder / 'Metadata'
        }
        
        ref_folders = {
            'video': ref_folder / 'Generated_Video',
            'metadata': ref_folder / 'Metadata'
        } if ref_folder and use_comparison else {}
        
        if not folders['metadata'].exists():
            logger.warning(f"Metadata folder not found: {folders['metadata']}")
            return []
        
        # OPTIMIZED: Use single-pass scanning for all folders
        reference_images, _, _ = self._scan_directory_once(folders['reference'])
        _, source_videos, _ = self._scan_directory_once(folders['source'])
        _, generated_videos, _ = self._scan_directory_once(folders['generated'])
        _, _, metadata_files = self._scan_directory_once(folders['metadata'])
        _, ref_videos, ref_metadata_raw = self._scan_directory_once(ref_folders.get('video', Path()))
        _, _, ref_metadata = self._scan_directory_once(ref_folders.get('metadata', Path()))
        
        # Batch load all metadata
        metadata_cache = self._load_json_batch(metadata_files) if metadata_files else {}
        ref_metadata_cache = self._load_json_batch(ref_metadata) if ref_metadata else {}
        
        logger.info(f"Runway files found: {len(reference_images)} refs, {len(source_videos)} sources, "
                   f"{len(generated_videos)} generated, {len(metadata_files)} metadata")
        
        # Pre-extract frames for all videos in parallel
        all_videos = list(source_videos.values()) + list(generated_videos.values())
        if all_videos:
            self._extract_frames_parallel(all_videos)
        
        pairs = []
        for stem, meta_path in metadata_files.items():
            md = metadata_cache.get(stem, {})
            if not md:
                continue
            
            # Find matching files using metadata references
            ref_img_path = next((p for p in reference_images.values()
                               if p.name == md.get('reference_image', '')), None)
            src_vid_path = next((p for p in source_videos.values()
                               if p.name == md.get('source_video', '')), None)
            gen_vid_path = next((p for p in generated_videos.values()
                               if p.name == md.get('generated_video', '')), None)
            
            # Determine source file and path
            source_file = None
            source_path = None
            if ref_img_path:
                source_file = ref_img_path.name
                source_path = ref_img_path
                logger.info("Face Swap task - using reference image as source")
            elif src_vid_path:
                source_file = src_vid_path.name
                source_path = src_vid_path
                src_vid_path = None  # Don't duplicate in source_video_path
                logger.info("Background Removal task - using source video as source")
            else:
                logger.warning(f"No valid source found for meta {stem}")
                continue
            
            # Handle reference comparisons
            ref_vid_path, ref_md = None, {}
            if use_comparison:
                r_base = stem.replace('runway_metadata', '')
                ref_vid_path = self.find_matching_video(r_base, ref_videos)
                ref_md = ref_metadata_cache.get(r_base, {})
            
            pair = MediaPair(
                source_file=source_file,
                source_path=source_path,
                api_type=self.api_name,
                generated_paths=[gen_vid_path] if gen_vid_path else [],
                reference_paths=[ref_vid_path] if ref_vid_path else [],
                source_video_path=src_vid_path,
                metadata=md,
                ref_metadata=ref_md,
                failed=not gen_vid_path or not md.get('success', False),
                ref_failed=use_comparison and (not ref_vid_path or not ref_md.get('success', False))
            )
            pairs.append(pair)
        
        logger.info(f"Created {len(pairs)} Runway media pairs")
        return pairs
    
    def create_wan_media_pairs(self, folder: Path, ref_folder: Optional[Path],
                              task: Dict, use_comparison: bool) -> List[MediaPair]:
        """
        Create Wan 2.2 media pairs.
        
        Structure:
        - Source Image/: Reference images
        - Source Video/: Source videos
        - Generated_Video/: Generated videos (named as image_video_mode.mp4)
        - Metadata/: JSON metadata files (named as image_video_metadata.json)
        """
        folders = {
            'source_image': folder / 'Source Image',
            'source_video': folder / 'Source Video',
            'generated': folder / 'Generated_Video',
            'metadata': folder / 'Metadata'
        }
        
        if not folders['metadata'].exists():
            logger.warning(f"Metadata folder not found: {folders['metadata']}")
            return []
        
        # OPTIMIZED: Use single-pass scanning for all folders
        source_images, _, _ = self._scan_directory_once(folders['source_image'])
        _, source_videos, _ = self._scan_directory_once(folders['source_video'])
        _, generated_videos, _ = self._scan_directory_once(folders['generated'])
        _, _, metadata_files = self._scan_directory_once(folders['metadata'])
        
        # Batch load all metadata
        metadata_cache = self._load_json_batch(metadata_files) if metadata_files else {}
        
        logger.info(f"Wan files found: {len(source_images)} images, {len(source_videos)} videos, "
                   f"{len(generated_videos)} generated, {len(metadata_files)} metadata")
        
        # Pre-extract frames for all videos in parallel
        all_videos = list(source_videos.values()) + list(generated_videos.values())
        if all_videos:
            self._extract_frames_parallel(all_videos)
        
        pairs = []
        for stem, meta_path in metadata_files.items():
            md = metadata_cache.get(stem, {})
            if not md:
                continue
            
            # Find matching files using metadata references
            src_img_path = next((p for p in source_images.values()
                               if p.name == md.get('source_image', '')), None)
            src_vid_path = next((p for p in source_videos.values()
                               if p.name == md.get('source_video', '')), None)
            gen_vid_path = next((p for p in generated_videos.values()
                               if p.name == md.get('generated_video', '')), None)
            
            if not src_img_path:
                logger.warning(f"No source image found for meta {stem}")
                continue
            
            # Use source video name as effect_name for per-video section dividers
            if src_vid_path:
                effect_name = src_vid_path.stem
            else:
                effect_name = self.config.get('effect') or self.config.get('effect_name')
                if not effect_name:
                    # Remove leading date (e.g., '1111 ') from folder name
                    m = re.match(r'^(\d{4})\s*(.+)', folder.name)
                    effect_name = m.group(2) if m else folder.name
            
            pair = MediaPair(
                source_file=src_img_path.name,
                source_path=src_img_path,
                api_type=self.api_name,
                generated_paths=[gen_vid_path] if gen_vid_path else [],
                reference_paths=[],
                source_video_path=src_vid_path,
                effect_name=effect_name,
                metadata=md,
                failed=not gen_vid_path or not md.get('success', False)
            )
            pairs.append(pair)
        
        logger.info(f"Created {len(pairs)} Wan media pairs")
        return pairs
    
    def process_base_folder_structure(self, task: Dict) -> List[MediaPair]:
        """Process base folder structure for vidu/pixverse/veo_itv APIs
        
        Args:
            task: If provided with 'effect' or 'folder' key, process only that single task.
                  Otherwise, process all tasks from config.
        """
        # veo_itv uses task-level folders, not base_folder
        if self.api_name == "veo_itv":
            logger.info(f"Processing {self.api_name} task folders")
            pairs = []
        else:
            base_folder = Path(self.config.get('base_folder', ''))
            if not base_folder.exists():
                logger.warning(f"Base folder not found: {base_folder}")
                return []
            
            logger.info(f"Processing {self.api_name} base folder: {base_folder}")
            pairs = []
        
        if self.api_name == "vidu_effects":
            # Process each effect folder
            # Support single-task filtering for grouped processing
            tasks_to_process = [task] if task.get('effect') else self.config.get('tasks', [])
            
            for task_config in tasks_to_process:
                effect = task_config.get('effect', '')
                category = task_config.get('category', 'Unknown')
                if not effect:
                    continue
                
                folders = {k: base_folder / effect / v
                          for k, v in {'src': 'Source', 'vid': 'Generated_Video', 'meta': 'Metadata'}.items()}
                
                if not folders['src'].exists():
                    logger.warning(f"Source folder not found for effect: {effect}")
                    continue
                
                logger.info(f"Processing Vidu effect: {effect}")
                
                # OPTIMIZED: Single-pass directory scanning
                images, _, _ = self._scan_directory_once(folders['src'])
                
                _, raw_videos, _ = self._scan_directory_once(folders['vid'])
                videos = {}
                for f in raw_videos.values():
                    key = self.extract_video_key(f.name, effect)
                    videos[key] = f
                
                _, _, metadata_files = self._scan_directory_once(folders['meta'])
                
                # Batch load metadata
                metadata_cache = self._load_json_batch(metadata_files) if metadata_files else {}
                
                logger.info(f"Images: {len(images)}, Videos: {len(videos)}, Meta {len(metadata_files)}")
                
                # Pre-compute aspect ratios
                all_media = list(images.values()) + list(videos.values())
                if all_media:
                    self._compute_aspect_ratios_batch(all_media, are_videos={p: True for p in videos.values()})
                
                # Match metadata to source files
                for key, img in images.items():
                    metadata = metadata_cache.get(key, {})
                    
                    vid = videos.get(key)
                    pair = MediaPair(
                        source_file=img.name,
                        source_path=img,
                        api_type=self.api_name,
                        generated_paths=[vid] if vid else [],
                        reference_paths=[],
                        effect_name=effect,
                        category=category,
                        metadata=metadata,
                        failed=not vid or not metadata.get('success', False)
                    )
                    pairs.append(pair)
        
        elif self.api_name == "vidu_reference":
            # Process vidu reference effects
            # Support single-task filtering for grouped processing
            if task.get('effect'):
                # Single effect mode (for grouping)
                effect_names = [task.get('effect')]
                logger.info(f"Processing single effect: {task.get('effect')}")
            else:
                # All effects mode (default)
                try:
                    effect_names = sorted([f.name for f in base_folder.iterdir()
                                         if f.is_dir() and not f.name.startswith('.') and (f / 'Source').exists()])
                    logger.info(f"Discovered {len(effect_names)} effect folders")
                except:
                    effect_names = [t.get('effect', '') for t in self.config.get('tasks', [])]
                    logger.info(f"Using {len(effect_names)} configured tasks")
            
            for effect in effect_names:
                if not effect:
                    continue
                
                folders = {k: base_folder / effect / v
                          for k, v in {'src': 'Source', 'vid': 'Generated_Video', 'meta': 'Metadata'}.items()}
                
                if not folders['src'].exists():
                    logger.warning(f"Source folder not found for effect: {effect}")
                    continue
                
                logger.info(f"Processing Vidu Reference effect: {effect}")
                
                # OPTIMIZED: Single-pass directory scanning
                images, _, _ = self._scan_directory_once(folders['src'])
                
                _, raw_videos, _ = self._scan_directory_once(folders['vid'])
                videos = {}
                for f in raw_videos.values():
                    videos[self.extract_key_reference(f.name, effect)] = f
                
                _, _, metadata_files = self._scan_directory_once(folders['meta'])
                
                # Batch load metadata
                metadata_cache = self._load_json_batch(metadata_files) if metadata_files else {}
                
                # Pre-compute aspect ratios
                all_media = list(images.values()) + list(videos.values())
                if all_media:
                    self._compute_aspect_ratios_batch(all_media, are_videos={p: True for p in videos.values()})
                
                # Scan Reference folder for reference images
                ref_folder = base_folder / effect / 'Reference'
                ref_images_map, _, _ = self._scan_directory_once(ref_folder) if ref_folder.exists() else ({}, {}, {})
                ref_image_list = sorted(ref_images_map.values(), key=lambda p: p.name.lower())
                
                # Create pairs
                for key, img in images.items():
                    metadata = metadata_cache.get(key, {})
                    
                    # Resolve reference images: prefer metadata list, fall back to folder scan
                    additional_source_paths = []
                    ref_names = metadata.get('reference_images', [])
                    if ref_names and ref_folder.exists():
                        for ref_name in ref_names:
                            ref_path = ref_folder / ref_name
                            if ref_path.exists():
                                additional_source_paths.append(ref_path)
                    if not additional_source_paths and ref_image_list:
                        additional_source_paths = list(ref_image_list)
                    
                    vid = videos.get(key)
                    pair = MediaPair(
                        source_file=img.name,
                        source_path=img,
                        api_type=self.api_name,
                        generated_paths=[vid] if vid else [],
                        reference_paths=[],
                        additional_source_paths=additional_source_paths,
                        effect_name=effect,
                        category="Reference",
                        metadata=metadata,
                        failed=not vid or not metadata.get('success', False)
                    )
                    pairs.append(pair)
        
        elif self.api_name == "pixverse":
            # Process pixverse effects
            # Support single-task filtering for grouped processing
            tasks_to_process = [task] if task.get('effect') else self.config.get('tasks', [])
            
            for task_config in tasks_to_process:
                effect = task_config.get('effect', '')
                category = task_config.get('category', 'Unknown')
                if not effect:
                    continue
                
                folders = {k: base_folder / effect / v
                          for k, v in {'src': 'Source', 'vid': 'Generated_Video', 'meta': 'Metadata'}.items()}
                
                if not folders['src'].exists():
                    continue
                
                # OPTIMIZED: Single-pass directory scanning
                images, _, _ = self._scan_directory_once(folders['src'])
                
                _, raw_videos, _ = self._scan_directory_once(folders['vid'])
                videos = {}
                for f in raw_videos.values():
                    key = self.extract_video_key(f.name, effect)
                    videos[key] = f
                
                _, _, metadata_files = self._scan_directory_once(folders['meta'])
                
                # Batch load metadata
                metadata_cache = self._load_json_batch(metadata_files) if metadata_files else {}
                
                # Pre-compute aspect ratios
                all_media = list(images.values()) + list(videos.values())
                if all_media:
                    self._compute_aspect_ratios_batch(all_media, are_videos={p: True for p in videos.values()})
                
                # Create pairs
                for key, img in images.items():
                    metadata = metadata_cache.get(key, {})
                    
                    vid = videos.get(key)
                    pair = MediaPair(
                        source_file=img.name,
                        source_path=img,
                        api_type=self.api_name,
                        generated_paths=[vid] if vid else [],
                        reference_paths=[],
                        effect_name=effect,
                        category=category,
                        metadata=metadata,
                        failed=not vid or not metadata.get('success', False)
                    )
                    pairs.append(pair)
        
        elif self.api_name == "kling_effects":
            # Process kling effects - uses custom_effect or effect as folder name
            # Support single-task filtering for grouped processing
            tasks_to_process = [task] if (task.get('effect') or task.get('custom_effect')) else self.config.get('tasks', [])
            
            for task_config in tasks_to_process:
                # custom_effect has priority over effect for folder name
                custom_effect = task_config.get('custom_effect', '')
                effect = custom_effect if custom_effect else task_config.get('effect', '')
                if not effect:
                    continue
                
                folders = {k: base_folder / effect / v
                          for k, v in {'src': 'Source', 'vid': 'Generated_Video', 'meta': 'Metadata'}.items()}
                
                if not folders['src'].exists():
                    logger.warning(f"Source folder not found for effect: {effect}")
                    continue
                
                logger.info(f"Processing Kling effect: {effect}")
                
                # OPTIMIZED: Single-pass directory scanning
                images, _, _ = self._scan_directory_once(folders['src'])
                
                _, raw_videos, _ = self._scan_directory_once(folders['vid'])
                videos = {}
                for f in raw_videos.values():
                    # Extract key by removing effect suffix pattern
                    key = self.extract_video_key(f.name, effect)
                    videos[key] = f
                
                _, _, metadata_files = self._scan_directory_once(folders['meta'])
                
                # Batch load metadata
                metadata_cache = self._load_json_batch(metadata_files) if metadata_files else {}
                
                logger.info(f"Images: {len(images)}, Videos: {len(videos)}, Meta: {len(metadata_files)}")
                
                # Pre-compute aspect ratios
                all_media = list(images.values()) + list(videos.values())
                if all_media:
                    self._compute_aspect_ratios_batch(all_media, are_videos={p: True for p in videos.values()})
                
                # Create pairs
                for key, img in images.items():
                    metadata = metadata_cache.get(key, {})
                    
                    vid = videos.get(key)
                    pair = MediaPair(
                        source_file=img.name,
                        source_path=img,
                        api_type=self.api_name,
                        generated_paths=[vid] if vid else [],
                        reference_paths=[],
                        effect_name=effect,
                        category="Effects",
                        metadata=metadata,
                        failed=not vid or not metadata.get('success', False)
                    )
                    pairs.append(pair)
        
        elif self.api_name == "veo_itv":
            # Process veo_itv - each task has its own folder with Source subfolder
            # Support single-task filtering for grouped processing
            tasks_to_process = [task] if task.get('folder') else self.config.get('tasks', [])
            
            for task_config in tasks_to_process:
                folder_path = Path(task_config.get('folder', ''))
                style_name = task_config.get('style_name', folder_path.name if folder_path else 'Unknown')
                if not folder_path or not folder_path.exists():
                    logger.warning(f"Folder not found: {folder_path}")
                    continue
                
                folders = {
                    'src': folder_path / 'Source',
                    'vid': folder_path / 'Generated_Video',
                    'meta': folder_path / 'Metadata'
                }
                
                if not folders['src'].exists():
                    logger.warning(f"Source folder not found: {folders['src']}")
                    continue
                
                logger.info(f"Processing Veo ITV style: {style_name}")
                
                # OPTIMIZED: Single-pass directory scanning
                images, _, _ = self._scan_directory_once(folders['src'])
                
                _, raw_videos, _ = self._scan_directory_once(folders['vid'])
                # Build video lookup: source_name -> {gen_num: video_path}
                videos_by_source = {}
                for f in raw_videos.values():
                    # Video naming: {source_name}_{gen_num}.mp4
                    stem = f.stem
                    if '_' in stem:
                        parts = stem.rsplit('_', 1)
                        source_key = self.normalize_key(parts[0])
                        try:
                            gen_num = int(parts[1])
                            if source_key not in videos_by_source:
                                videos_by_source[source_key] = {}
                            videos_by_source[source_key][gen_num] = f
                        except ValueError:
                            # Not a numbered video, use full stem as key
                            key = self.normalize_key(stem)
                            if key not in videos_by_source:
                                videos_by_source[key] = {}
                            videos_by_source[key][1] = f
                
                _, _, metadata_files = self._scan_directory_once(folders['meta'])
                
                # Batch load metadata - keyed by full filename (source_name_gen_num)
                metadata_cache = self._load_json_batch(metadata_files) if metadata_files else {}
                
                logger.info(f"Images: {len(images)}, Videos: {len(raw_videos)}, Meta: {len(metadata_files)}")
                
                # Pre-compute aspect ratios
                all_media = list(images.values()) + list(raw_videos.values())
                if all_media:
                    self._compute_aspect_ratios_batch(all_media, are_videos={p: True for p in raw_videos.values()})
                
                # Create pairs for each source image and its generated videos
                for key, img in images.items():
                    source_videos = videos_by_source.get(key, {})
                    
                    if source_videos:
                        # Create a pair for each generated video
                        for gen_num, vid in sorted(source_videos.items()):
                            meta_key = f"{key}_{gen_num}"
                            metadata = metadata_cache.get(meta_key, {})
                            
                            pair = MediaPair(
                                source_file=img.name,
                                source_path=img,
                                api_type=self.api_name,
                                generated_paths=[vid],
                                reference_paths=[],
                                effect_name=style_name,
                                category="ITV",
                                metadata=metadata,
                                failed=not metadata.get('success', False)
                            )
                            pairs.append(pair)
                    else:
                        # No videos generated yet
                        pair = MediaPair(
                            source_file=img.name,
                            source_path=img,
                            api_type=self.api_name,
                            generated_paths=[],
                            reference_paths=[],
                            effect_name=style_name,
                            category="ITV",
                            metadata={},
                            failed=True
                        )
                        pairs.append(pair)
        
        return pairs
    
    def process_genvideo_batch(self, task: Dict) -> List[MediaPair]:
        """Process GenVideo batch"""
        folder = Path(task['folder'])
        source_folder = folder / 'Source'
        generated_folder = folder / 'Generated_Image'
        metadata_folder = folder / 'Metadata'
        pairs = []
        
        if not source_folder.exists():
            logger.warning(f"Source folder not found: {source_folder}")
            return pairs
        
        if not generated_folder.exists():
            logger.warning(f"Generated folder not found: {generated_folder}")
            return pairs
        
        # Process each source image
        source_images = [f for f in source_folder.iterdir()
                        if f.suffix.lower() in self.IMAGE_EXTS]
        
        # Pre-load all possible metadata files
        potential_metadata = {}
        if metadata_folder.exists():
            for mf in metadata_folder.iterdir():
                if mf.suffix.lower() == '.json':
                    try:
                        with open(mf, 'r', encoding='utf-8') as f:
                            potential_metadata[mf.stem] = json.load(f)
                    except Exception as e:
                        logger.warning(f"Failed to load metadata {mf.name}: {e}")
        
        for src_img in source_images:
            basename = src_img.stem
            
            # Try to find metadata from pre-loaded cache
            metadata = {}
            for key_pattern in [f"{basename}_{src_img.name}_metadata", f"{basename}_metadata", basename]:
                if key_pattern in potential_metadata:
                    metadata = potential_metadata[key_pattern]
                    logger.info(f"Found metadata for {basename}")
                    break
            
            # Find generated image
            gen_img = generated_folder / f"{basename}.jpg"
            if not gen_img.exists():
                for ext in ['.png', '.jpeg', '.webp']:
                    alt_gen = generated_folder / f"{basename}_generated{ext}"
                    if alt_gen.exists():
                        gen_img = alt_gen
                        break
            
            pair = MediaPair(
                source_file=src_img.name,
                source_path=src_img,
                api_type='genvideo',
                generated_paths=[gen_img] if gen_img.exists() else [],
                reference_paths=[],
                metadata=metadata,
                failed=not gen_img.exists() or not metadata.get('success', False)
            )
            pairs.append(pair)
            
            if pair.failed:
                logger.warning(f"Failed pair: {src_img.name}")
            else:
                logger.info(f"Valid pair: {src_img.name} -> {gen_img.name}")
        
        logger.info(f"Created {len(pairs)} GenVideo media pairs")
        return pairs
    
    def process_text_to_video_batch(self, task: Dict) -> List[MediaPair]:
        """Process text-to-video APIs (Veo, Kling TTV)"""
        # Get root folder from config for kling_ttv, or task-level for veo
        if self.api_name == 'kling_ttv':
            root_folder = Path(self.config.get('output_folder', task.get('output_folder', '')))
            output_folder = root_folder / 'Generated_Video'
            metadata_folder = root_folder / 'Metadata'
        else:  # veo
            # For veo, output_folder in config already points to Generated_Video directory
            output_folder = Path(task.get('output_folder', ''))
            metadata_folder = output_folder.parent / 'Metadata'
        
        pairs = []
        
        if not output_folder.exists():
            logger.warning(f"Output folder not found: {output_folder}")
            return pairs
        
        # Get generated videos and metadata
        _, generated_videos, _ = self._scan_directory_once(output_folder)
        _, _, metadata_files = self._scan_directory_once(metadata_folder)
        
        # Batch load metadata
        metadata_cache = self._load_json_batch(metadata_files) if metadata_files else {}
        
        logger.info(f"Text-to-video files found: {len(generated_videos)} generated, {len(metadata_files)} metadata")
        
        # Pre-extract frames for all videos
        if generated_videos:
            self._extract_frames_parallel(list(generated_videos.values()))
        
        # Create pairs from metadata (metadata drives the pairing for text-to-video)
        for stem, meta_path in metadata_files.items():
            md = metadata_cache.get(stem, {})
            if not md:
                continue
            
            # Find matching generated video
            gen_vid_path = next((p for p in generated_videos.values()
                               if p.name == md.get('generated_video', '')), None)
            
            # Get style name for display
            style_name = md.get('style_name', 'Unknown')
            gen_num = md.get('generation_number', 1)
            
            # For text-to-video, source is the prompt (no source file)
            pair = MediaPair(
                source_file=f"{style_name}-{gen_num}",
                source_path=None,  # No source file for text-to-video
                api_type=self.api_name,
                generated_paths=[gen_vid_path] if gen_vid_path else [],
                reference_paths=[],
                effect_name=style_name,
                metadata=md,
                failed=not gen_vid_path or not md.get('success', False)
            )
            pairs.append(pair)
        
        logger.info(f"Created {len(pairs)} {self.api_name} media pairs")
        return pairs
    
    # ================== UTILITY METHODS ==================
    
    def load_config(self):
        """Load API-specific configuration from YAML or JSON"""
        # Skip loading if config was already set programmatically
        if self.config:
            logger.info("✓ Using pre-set configuration (runtime overrides applied)")
            return
        
        config_path = Path(self.config_file)
        try:
            with open(config_path, 'r', encoding='utf-8') as f:
                # Detect format by extension
                if config_path.suffix.lower() in ['.yaml', '.yml']:
                    self.config = yaml.safe_load(f)
                    logger.info(f"✓ Config loaded: {self.config_file} (YAML)")
                else:
                    self.config = json.load(f)
                    logger.info(f"✓ Config loaded: {self.config_file} (JSON)")
        except Exception as e:
            logger.error(f"✗ Config error: {e}")
            sys.exit(1)

    def set_config(self, config: dict) -> None:
        """
        Set configuration directly, bypassing file loading.
        
        This method allows the GUI and programmatic callers to inject
        a pre-merged configuration dictionary (with runtime overrides applied)
        without reading from a file.
        
        Args:
            config: Configuration dictionary to use for report generation.
        """
        self.config = config
        
        # Update Kling display name if applicable
        if self.api_name in ['kling', 'kling_endframe', 'kling_ttv', 'kling_motion']:
            self._update_kling_display_name()
    
    def _update_kling_display_name(self):
        """Update Kling display name based on model in config"""
        # Check both 'model' and 'model_version' fields at root level
        model = (self.config.get('model') or self.config.get('model_version', '')).lower()
        
        # For kling_ttv/kling_motion, check the first task's model if not at root level
        if not model and self.api_name in ('kling_ttv', 'kling_motion'):
            tasks = self.config.get('tasks', [])
            if tasks and 'model' in tasks[0]:
                model = tasks[0]['model'].lower()
        
        # Map official Kling model names to display names
        # Official model names: v1.5, v1.6, v2.0-master, v2.1, v2.1-master, v2.5-turbo, v2.6, v3
        model_mapping = {
            'v1.5': 'Kling 1.5',
            'v1.6': 'Kling 1.6',
            'v2.0-master': 'Kling 2.0',
            'v2.1': 'Kling 2.1',
            'v2.1-master': 'Kling 2.1',
            'v2.5-turbo': 'Kling 2.5',
            'v2.6': 'Kling 2.6',
            'v3': 'Kling 3.0',
        }
        
        # Try to find a match
        display_name = model_mapping.get(model)
        
        # If no exact match, try to extract version numbers
        if not display_name and model:
            # Try to extract version like "v2.1" or "2.1" from the model string
            import re
            match = re.search(r'(?:v)?(\d+[._-]\d+)(?:[._-]?turbo)?', model)
            if match:
                version = match.group(1).replace('_', '.').replace('-', '.')
                is_turbo = 'turbo' in model
                display_name = f"Kling {version}{' Turbo' if is_turbo else ''}"
        
        # Update the display name for kling, kling_endframe, kling_ttv, and kling_motion
        if display_name:
            if self.api_name == 'kling_endframe':
                self._api_display_names['kling_endframe'] = f"{display_name} Endframe"
            elif self.api_name == 'kling_ttv':
                self._api_display_names['kling_ttv'] = f"{display_name} TTV"
            elif self.api_name == 'kling_motion':
                self._api_display_names['kling_motion'] = f"{display_name} Motion"
            else:
                self._api_display_names['kling'] = display_name
            logger.info(f"✓ Kling display name set to: {self._api_display_names[self.api_name]}")
        else:
            # Default values
            if self.api_name == 'kling_endframe':
                self._api_display_names['kling_endframe'] = 'Kling 1.6 Endframe'  # Default for endframe
            elif self.api_name == 'kling_ttv':
                self._api_display_names['kling_ttv'] = 'Kling 1.6 TTV'  # Default for TTV
            elif self.api_name == 'kling_motion':
                self._api_display_names['kling_motion'] = 'Kling 3.0 Motion'  # Default for motion
            else:
                self._api_display_names['kling'] = 'Kling 2.1'  # Default for regular kling
            if model:
                logger.warning(f"⚠ Could not determine Kling version from model '{model}', defaulting to '{self._api_display_names[self.api_name]}'")
            else:
                logger.info(f"ℹ No model version specified in config, defaulting to '{self._api_display_names[self.api_name]}'")
    
    def load_report_definitions(self):
        """Load report definitions from api_definitions.json"""
        # First try using the path helper (works for PyInstaller bundles)
        api_def_path = get_core_path("api_definitions.json")
        
        if api_def_path.exists():
            try:
                with open(api_def_path, 'r', encoding='utf-8') as f:
                    all_definitions = json.load(f)
                self.report_definitions = all_definitions.get(self.api_name, {}).get('report', {})
                logger.info(f"✓ API definitions loaded from: {api_def_path}")
                return
            except Exception as e:
                logger.warning(f"⚠ Error loading API definitions: {e}")
        
        # Fallback to relative paths
        definition_paths = [
            "core/api_definitions.json",
            "api_definitions.json",
            "config/api_definitions.json"
        ]
        
        for def_path in definition_paths:
            try:
                with open(def_path, 'r', encoding='utf-8') as f:
                    all_definitions = json.load(f)
                self.report_definitions = all_definitions.get(self.api_name, {}).get('report', {})
                logger.info(f"✓ API definitions loaded from: {def_path}")
                return
            except Exception:
                continue
        
        logger.warning(f"⚠ API definitions not found, using defaults")
        self.set_default_report_definitions()
    
    def set_default_report_definitions(self):
        """Set default report definitions"""
        self.report_definitions = {
            "enabled": True,
            "template_path": "templates/I2V templates.pptx",
            "comparison_template_path": "templates/I2V Comparison Template.pptx",
            "output_directory": "/Users/ethanhsu/Desktop/EthanHsu-cl/GAI/Report",
            "use_comparison": self.api_name in ["kling", "nano_banana", "runway", "wan", "dreamactor", "kling_motion"]
        }
    
    def _extract_date_from_folder(self, folder):
        """Extract date from folder name or use current date"""
        folder_name = Path(folder).name if isinstance(folder, (str, Path)) else str(folder)
        m = re.match(r'(\d{4})\s*(.+)', folder_name)
        return m.group(1) if m else datetime.now().strftime("%m%d")
    
    def get_cmp_filename(self, folder1: str, folder2: str, model: str = '', effect_names1=None, effect_names2=None) -> tuple:
        """Generate comparison filename using API name and effect names.
        
        Returns:
            tuple: (api_line, styles_line) where api_line contains date and model,
                   and styles_line contains the effect names comparison.
        """
        d = self._extract_date_from_folder(folder1)
        
        # Use model (API name) as the primary identifier
        # Effect names are the actual content description
        effect_str1 = ', '.join(effect_names1) if effect_names1 else 'Test'
        effect_str2 = ', '.join(effect_names2) if effect_names2 else 'Reference'
        
        # Build API line (date + model)
        api_parts = [f"[{d}]"]
        if model:
            api_parts.append(model)
        api_line = ' '.join(api_parts)
        
        # Build styles line
        styles_line = f"{effect_str1} vs {effect_str2}"
        
        return (api_line, styles_line)

    def get_filename(self, folder, model='', effect_names=None):
        """Generate filename using API name and effect names.
        
        Returns:
            tuple: (api_line, styles_line) where api_line contains date and model,
                   and styles_line contains the effect/style names.
        """
        # Handle grouped tasks
        if isinstance(folder, dict) and folder.get('_is_grouped'):
            return self._get_grouped_filename(folder, model, effect_names)
        
        d = self._extract_date_from_folder(folder)
        
        # Use model (API name) as the primary identifier
        # For APIs with many styles, show count instead of listing all names to avoid long filenames
        if self.api_name in ['veo', 'kling_ttv', 'veo_itv', 'wan', 'dreamactor', 'kling_motion'] and effect_names:
            effect_str = f"{len(effect_names)} {'Style' if len(effect_names) == 1 else 'Styles'}"
        else:
            # Effect names are the actual content description
            effect_str = ', '.join(effect_names) if effect_names else 'Test'
            # Only truncate very long effect strings (not short concatenations)
            if len(effect_str) > 60:
                effect_str = effect_str[:60] + '...'
        
        # Build API line (date + model)
        api_parts = [f"[{d}]"]
        if model:
            api_parts.append(model)
        api_line = ' '.join(api_parts)
        
        return (api_line, effect_str)
    
    def _get_grouped_filename(self, grouped_task: Dict, model: str = '', effect_names=None) -> tuple:
        """Generate filename for grouped tasks.
        
        Handles both folder-based and base-folder APIs.
        
        Returns:
            tuple: (api_line, styles_line) where api_line contains date and model,
                   and styles_line contains the effect/style names.
        """
        is_base_folder_api = grouped_task.get('_is_base_folder_api', False)
        group_num = grouped_task.get('_group_number', 1)
        
        if is_base_folder_api:
            # Base folder API (vidu_effects, etc.) - use base folder for date and effect names
            base_folder = grouped_task.get('base_folder', '')
            d = self._extract_date_from_folder(base_folder) if base_folder else datetime.now().strftime("%m%d")
            
            # Use effect names from the grouped task
            effect_list = grouped_task.get('_effect_names', [])
            # For APIs with many styles, show count instead of listing all names
            if self.api_name in ['veo', 'kling_ttv', 'veo_itv'] and effect_list:
                effect_str = f"{len(effect_list)} {'Style' if len(effect_list) == 1 else 'Styles'}"
            else:
                effect_str = ', '.join(effect_list) if effect_list else 'Combined Effects'
        else:
            # Folder-based API (nano_banana, kling, etc.) - use folder names
            folder_names = grouped_task.get('_folder_names', [])
            
            # Extract date from first folder (prioritize folder date over current date)
            # For veo_itv, use parent folder name (contains date like "0130 6 Styles")
            if self.api_name == "veo_itv":
                parent_folder = grouped_task.get('_parent_folder_name', '')
                d = self._extract_date_from_folder(parent_folder) if parent_folder else datetime.now().strftime("%m%d")
            else:
                d = self._extract_date_from_folder(folder_names[0]) if folder_names else datetime.now().strftime("%m%d")
            
            # Build effect string - combine all unique effects
            # For APIs with many styles, show count instead of listing all names
            if self.api_name in ['veo', 'kling_ttv', 'veo_itv', 'wan', 'dreamactor', 'kling_motion'] and effect_names:
                effect_str = f"{len(effect_names)} {'Style' if len(effect_names) == 1 else 'Styles'}"
            else:
                effect_str = ', '.join(effect_names) if effect_names else 'Combined'
        
        # Build API line (date + model)
        api_parts = [f"[{d}]"]
        if model:
            api_parts.append(model)
        api_line = ' '.join(api_parts)
        
        return (api_line, effect_str)

    
    def _scan_directory_once(self, folder: Path, image_exts=None, video_exts=None, metadata_exts=None):
        """Scan directory once and categorize files by type - major performance optimization"""
        image_exts = image_exts or self.IMAGE_EXTS
        video_exts = video_exts or self.VIDEO_EXTS
        metadata_exts = metadata_exts or self.METADATA_EXTS
        
        images, videos, metadata = {}, {}, {}
        
        if not folder or not folder.exists():
            return images, videos, metadata
        
        for f in folder.iterdir():
            if not f.is_file():
                continue
            suffix = f.suffix.lower()
            if suffix in image_exts:
                images[self.normalize_key(f.stem)] = f
            elif suffix in video_exts:
                videos[self.normalize_key(f.stem)] = f
            elif suffix in metadata_exts:
                key = f.stem.replace('_metadata', '')
                metadata[self.normalize_key(key)] = f
        
        return images, videos, metadata
    
    def find_matching_video(self, base_name: str, video_files: dict) -> Optional[Path]:
        """Enhanced video matching for all APIs"""
        if base_name in video_files:
            return video_files[base_name]
        
        for v_name, v_path in video_files.items():
            if v_name.startswith(base_name):
                return v_path
        
        return None
    
    def get_aspect_ratio(self, path, is_video=False):
        """Calculate aspect ratio with caching - always use actual dimensions"""
        key = str(path)
        if key in self._ar_cache: 
            return self._ar_cache[key]
        
        # Try to get actual dimensions from the file first
        try:
            if is_video and cv2:
                cap = cv2.VideoCapture(str(path))
                if cap.isOpened():
                    w, h = cap.get(3), cap.get(4)
                    cap.release()
                    if w > 0 and h > 0:
                        ar = w / h
                        self._ar_cache[key] = ar
                        return ar
            else:
                with Image.open(path) as img:
                    ar = img.width / img.height
                    self._ar_cache[key] = ar
                    return ar
        except:
            pass
        
        # Fallback: Use filename patterns only if we can't read the file
        import re
        fn = path.name.lower()
        if re.search(r'(?:^|_|-|\s)9[_-]16(?:$|_|-|\s)', fn) or 'portrait' in fn: 
            return 9/16
        if re.search(r'(?:^|_|-|\s)1[_-]1(?:$|_|-|\s)', fn) or 'square' in fn: 
            return 1
        if re.search(r'(?:^|_|-|\s)16[_-]9(?:$|_|-|\s)', fn) or 'landscape' in fn: 
            return 16/9
        
        # Final fallback
        return 16/9
    
    def extract_first_frame(self, video_path):
        """Extract first frame with caching"""
        if not cv2:
            return None
        
        video_key = str(video_path)
        if video_key in self._frame_cache:
            return self._frame_cache[video_key]
        
        try:
            cap = cv2.VideoCapture(str(video_path))
            if not cap.isOpened():
                return None
            
            ret, frame = cap.read()
            cap.release()
            
            if not ret or frame is None:
                return None
            
            # Create temporary file for the frame
            temp_dir = tempfile.gettempdir()
            frame_filename = f"frame_{video_path.stem}_{hash(video_key) % 10000}.jpg"
            frame_path = Path(temp_dir) / frame_filename
            
            # Save frame as JPEG
            cv2.imwrite(str(frame_path), frame)
            
            # Cache the result
            self._frame_cache[video_key] = str(frame_path)
            return str(frame_path)
        except Exception as e:
            logger.warning(f"Failed to extract frame from {video_path}: {e}")
            return None
    
    def _extract_frames_parallel(self, video_paths):
        """Extract first frames from multiple videos in parallel - major speedup"""
        if not cv2 or not video_paths:
            return {}
        
        def extract_one(video_path):
            return (video_path, self.extract_first_frame(video_path))
        
        paths_list = list(video_paths) if not isinstance(video_paths, list) else video_paths
        
        with ThreadPoolExecutor(max_workers=self._max_workers) as executor:
            if self._show_progress and len(paths_list) > 10:
                pbar = tqdm(total=len(paths_list), desc="Extracting video frames", unit="videos")
                results = []
                for result in executor.map(lambda p: extract_one(p), paths_list):
                    results.append(result)
                    pbar.update()
                pbar.close()
            else:
                results = list(executor.map(lambda p: extract_one(p), paths_list))
        
        return dict(results)
    
    def _compute_aspect_ratios_batch(self, media_paths, are_videos=False):
        """Pre-compute aspect ratios for multiple files in parallel"""
        if not media_paths:
            return
        
        def compute_one(path):
            ar = self.get_aspect_ratio(path, are_videos.get(path, False) if isinstance(are_videos, dict) else are_videos)
            return (str(path), ar)
        
        paths_list = list(media_paths) if not isinstance(media_paths, list) else media_paths
        
        with ThreadPoolExecutor(max_workers=self._max_workers) as executor:
            if self._show_progress and len(paths_list) > 20:
                pbar = tqdm(total=len(paths_list), desc="Computing aspect ratios", unit="files")
                results = []
                for result in executor.map(compute_one, paths_list):
                    results.append(result)
                    pbar.update()
                pbar.close()
            else:
                results = list(executor.map(compute_one, paths_list))
        
        # Update cache
        for path_str, ar in results:
            self._ar_cache[path_str] = ar
    
    def _process_in_batches(self, items, process_func, batch_size=None, desc="Processing"):
        """Process items in optimal batches to manage memory usage"""
        if not items:
            return []
        
        batch_size = batch_size or self._batch_size
        results = []
        
        total_batches = (len(items) + batch_size - 1) // batch_size
        
        if self._show_progress and len(items) > batch_size:
            logger.info(f"Processing {len(items)} items in {total_batches} batches of {batch_size}")
        
        for i in range(0, len(items), batch_size):
            batch = items[i:i + batch_size]
            batch_results = process_func(batch)
            results.extend(batch_results)
            
            # Clear temporary caches after each batch to manage memory
            if i + batch_size < len(items):
                import gc
                gc.collect()
        
        return results
    
    def extract_video_key(self, filename: str, effect_name: str) -> str:
        """Extract video key - FIXED for proper effect name removal (memoized)"""
        cache_key = f"{filename}|{effect_name}"
        if cache_key in self._extract_key_cache:
            return self._extract_key_cache[cache_key]
        
        stem = Path(filename).stem
        
        # First remove the _effect suffix
        stem = re.sub(r"_effect$", "", stem, flags=re.IGNORECASE)
        
        # Create multiple possible effect patterns to match
        effect_variations = [
            effect_name.replace(' ', '_'),  # Space to underscore
            effect_name.replace('-', '_'),  # Dash to underscore
            effect_name.replace(' ', '_').replace('-', '_'),  # Both replacements
            effect_name  # Original with spaces/dashes
        ]
        
        # Try to remove each variation (case insensitive)
        for effect_var in effect_variations:
            # Remove with underscore prefix
            pattern = f"_{re.escape(effect_var)}"
            stem = re.sub(pattern, "", stem, flags=re.IGNORECASE)
            
            # Remove without underscore prefix
            pattern = re.escape(effect_var)
            stem = re.sub(pattern, "", stem, flags=re.IGNORECASE)
        
        # Clean up any trailing underscores or effect patterns
        for pattern in [r"_generated", r"_output", r"_result"]:
            stem = re.sub(pattern, "", stem, flags=re.IGNORECASE)
        
        result = self.normalize_key(stem)
        self._extract_key_cache[cache_key] = result
        return result
    
    def extract_key_reference(self, filename: str, effect: str) -> str:
        """Extract key for vidu_reference - handles effect name with spaces/dashes/underscores (memoized)"""
        cache_key = f"{filename}|{effect}"
        if cache_key in self._extract_key_cache:
            return self._extract_key_cache[cache_key]
        
        stem = Path(filename).stem
        
        # Create all possible variations of the effect name as it might appear in filename
        effect_variations = [
            effect.replace(' ', '_').replace('-', '_'),  # "Corpse_Bride_V3_2"
            effect.replace(' ', '_'),  # Keep dashes: "Corpse_Bride_V3-2"
            effect.replace('-', '_'),  # "Corpse Bride V3_2"
            effect,  # Original
        ]
        
        # Try removing each variation from the end of the stem
        for eff_var in effect_variations:
            eff_lower = eff_var.lower()
            stem_lower = stem.lower()
            
            # Try with underscore separator
            if stem_lower.endswith(f'_{eff_lower}'):
                stem = stem[:-(len(eff_lower) + 1)]
                break
            # Try without separator (adjacent)
            if stem_lower.endswith(eff_lower):
                stem = stem[:-len(eff_lower)]
                break
        
        # Clean up any trailing underscores
        stem = stem.rstrip('_')
        
        result = self.normalize_key(stem)
        self._extract_key_cache[cache_key] = result
        return result
    
    def cleanup_temp_frames(self):
        """Clean up temporary frame files"""
        for frame_path in self._frame_cache.values():
            try:
                if Path(frame_path).exists():
                    os.unlink(frame_path)
            except Exception:
                pass
        self._frame_cache.clear()
    
    def cleanup_tempfiles(self):
        """Clean up temporary files from image format conversions"""
        for f in self._tempfiles_to_cleanup:
            try:
                if Path(f).exists():
                    os.unlink(f)
            except Exception:
                pass
        self._tempfiles_to_cleanup.clear()
    
    def cleanup_caches(self):
        """Clear all memory caches - useful between large batches"""
        self._normalize_cache.clear()
        self._extract_key_cache.clear()
        logger.debug(f"Cleared string caches: {len(self._normalize_cache)} + {len(self._extract_key_cache)} entries")
    
    def configure_performance(self, batch_size=None, max_workers=None, show_progress=None):
        """Configure performance settings for optimization"""
        if batch_size is not None:
            self._batch_size = batch_size
        if max_workers is not None:
            self._max_workers = max_workers
        if show_progress is not None:
            self._show_progress = show_progress and HAS_TQDM
        
        logger.info(f"Performance config: batch_size={self._batch_size}, "
                   f"max_workers={self._max_workers}, progress={self._show_progress}")
    
    # ================== PRESENTATION CREATION ==================
    
    def _gather_effect_names(self, pairs: List[MediaPair]) -> List[str]:
        """Extract unique effect names from pairs in order of appearance"""
        effect_names = []
        seen = set()
        for p in pairs:
            ename = p.effect_name.strip() if hasattr(p, 'effect_name') and p.effect_name else None
            if ename and ename not in seen:
                effect_names.append(ename)
                seen.add(ename)
        return effect_names
    
    def _load_presentation_template(self, task: Dict) -> tuple:
        """Load presentation template and return (ppt, template_loaded, use_comparison)"""
        # Stacked-layout APIs place media at custom positions, so they always
        # use the standard 2-media template even when use_comparison_template is set.
        slide_config = self.get_slide_config()
        if slide_config.get('override_positions', False):
            use_comparison = False
        else:
            use_comparison = task.get('use_comparison_template', False) or bool(task.get('reference_folder'))
        template_key = 'comparison_template_path' if use_comparison else 'template_path'
        template_path = (self.config.get(template_key) or
                        self.report_definitions.get(template_key,
                        'templates/I2V Comparison Template.pptx' if use_comparison else 'templates/I2V templates.pptx'))
        
        # Resolve template path - check if it exists, if not try using resource path
        template_path_obj = Path(template_path)
        if not template_path_obj.exists():
            # Try using the app base path for bundled resources
            resource_template = get_resource_path(template_path)
            if resource_template.exists():
                template_path = str(resource_template)
        
        try:
            ppt = Presentation(template_path) if Path(template_path).exists() else Presentation()
            template_loaded = Path(template_path).exists()
            logger.info(f"✓ Template loaded: {template_path}")
        except Exception as e:
            logger.warning(f"⚠ Template load failed: {e}, using blank presentation")
            ppt = Presentation()
            template_loaded = False
        
        # Set slide dimensions
        ppt.slide_width, ppt.slide_height = Cm(33.87), Cm(19.05)
        
        return ppt, template_loaded, use_comparison
    
    def create_presentation(self, pairs: List[MediaPair], task: Dict) -> bool:
        """Create presentation using unified system"""
        if not pairs:
            logger.warning("No media pairs to process")
            return False

        # Gather effect names from pairs
        effect_names = self._gather_effect_names(pairs)
        
        # Load template and get configuration
        ppt, template_loaded, use_comparison = self._load_presentation_template(task)

        # Create title slide
        self.create_title_slide(ppt, task, use_comparison, effect_names)

        # Create content slides using UNIFIED SYSTEM
        self.create_slides(ppt, pairs, template_loaded, use_comparison)

        # Save presentation
        return self.save_presentation(ppt, task, use_comparison, effect_names)
    
    def create_grouped_presentation(self, task_pairs_list: List[Dict], combined_task: Dict) -> bool:
        """Create a grouped presentation with individual title slides for each task"""
        if not task_pairs_list:
            logger.warning("No task pairs to process")
            return False
        
        # Gather all effect names from all tasks for the overall filename
        all_pairs = [p for item in task_pairs_list for p in item['pairs']]
        all_effect_names = self._gather_effect_names(all_pairs)
        
        # Load template using first task's settings
        first_task = task_pairs_list[0]['task']
        ppt, template_loaded, use_comparison = self._load_presentation_template(first_task)
        
        # Create overall title slide with all effect names (only once for the whole presentation)
        self.create_title_slide(ppt, combined_task, use_comparison, all_effect_names)
        
        # Process each task individually WITHOUT creating individual title slides
        # (We only want ONE title slide at the top with all effect names)
        for idx, item in enumerate(task_pairs_list, 1):
            task = item['task']
            pairs = item['pairs']
            
            # Get effect names for this specific task
            task_effect_names = self._gather_effect_names(pairs)
            
            # NOTE: Do NOT create individual title slides for each task in grouped mode
            # The overall title slide at the top already shows all effects
            
            # Create content slides for this task
            slide_config = self.get_slide_config()
            if slide_config.get('override_positions', False):
                task_use_comparison = False
            else:
                task_use_comparison = task.get('use_comparison_template', False) or bool(task.get('reference_folder'))
            self.create_slides(ppt, pairs, template_loaded, task_use_comparison)
            
            logger.info(f"  ✓ Added task {idx}/{len(task_pairs_list)}: {len(pairs)} slides")
        
        # Save with combined filename
        return self.save_presentation(ppt, combined_task, use_comparison, all_effect_names)

    def create_title_slide(self, ppt: Presentation, task: Dict, use_comparison: bool, effect_names=None):
        """Create title slide"""
        if not ppt.slides:
            return

        # Get folder names for title generation
        if task.get('_is_grouped'):
            # For grouped tasks, pass the entire task dict as folder_name
            folder_name = task
        elif self.api_name in ["vidu_effects", "vidu_reference", "pixverse"]:
            folder_name = Path(self.config.get('base_folder', '')).name
        elif self.api_name == "veo_itv":
            # For veo_itv, get parent folder (e.g., "0130 6 Styles" from "0130 6 Styles/Street Rap")
            # since the date prefix is in the parent, not the style folder
            folder_path = Path(task.get('folder', ''))
            folder_name = folder_path.parent.name
        else:
            folder_name = task.get('folder', Path(self.config.get('base_folder', '')).name)
            if isinstance(folder_name, str):
                folder_name = Path(folder_name).name

        api_display = self._api_display_names.get(self.api_name, self.api_name.title())

        # Generate title (now returns tuple of api_line, styles_line)
        if use_comparison and task.get('reference_folder'):
            ref_name = Path(task['reference_folder']).name
            api_line, styles_line = self.get_cmp_filename(folder_name, ref_name, api_display, effect_names1=effect_names)
        else:
            api_line, styles_line = self.get_filename(folder_name, api_display, effect_names=effect_names)

        # Update title slide with two-line formatting
        if ppt.slides and ppt.slides[0].shapes:
            title_shape = ppt.slides[0].shapes[0]
            tf = title_shape.text_frame
            tf.clear()
            
            # First paragraph: API line (uses default/existing font size)
            p1 = tf.paragraphs[0]
            p1.text = api_line
            p1.alignment = PP_ALIGN.CENTER
            
            # Second paragraph: Styles line with font size 36
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            run = p2.add_run()
            run.text = styles_line
            run.font.size = Pt(36)

        # Add links
        self.add_links(ppt, task)

    def add_links(self, ppt: Presentation, task: Dict):
        """Add hyperlinks to title slide"""
        if not ppt.slides:
            return

        slide = ppt.slides[0]

        # Find or create info box (skip the title shape which is shapes[0])
        # Only search shapes[1:] to avoid accidentally selecting the title shape
        info_box = next((s for s in list(slide.shapes)[1:] if hasattr(s,'text_frame') and s.text_frame.text and 
                        any(k in s.text_frame.text.lower() for k in ['design','testbed','source'])), None)

        if not info_box:
            info_box = slide.shapes.add_textbox(Cm(5), Cm(13), Cm(20), Cm(4))

        info_box.text_frame.clear()

        # Get API-specific links
        testbed_url = self.config.get('testbed', f'http://192.168.31.161/external-testbed/{self.api_name}/')
        
        # Get design link - combine root + task-level if both exist
        root_design_link = self.config.get('root_design_link', '')
        task_design_link = task.get('design_link', '')
        
        if root_design_link and task_design_link:
            # Combine root link with task-level anchor/suffix
            design_link = root_design_link + task_design_link
        elif root_design_link:
            design_link = root_design_link
        else:
            design_link = task_design_link
        
        # Check if this is a grouped presentation
        is_grouped = task.get('_is_grouped', False)
        is_base_folder_api = task.get('_is_base_folder_api', False)
        
        links = []
        
        # Add design link (same for all)
        if design_link:
            links.append(("Design: ", "Link", design_link))
        
        # Add testbed link
        links.append(("Testbed: ", testbed_url, testbed_url))
        
        # Handle source video links
        if is_grouped and not is_base_folder_api:
            # Folder-based APIs (nano_banana, kling, runway, genvideo) with grouping
            # Check if there's a root source+video link in config
            root_source_link = self.config.get('root_source_video_link', '')
            
            if root_source_link:
                # If root link exists, add it as "Link" with hyperlink and ignore task-specific links
                links.append(("Source + Video: ", "Link", root_source_link))
            else:
                # No root link - add "Source + Video: " without hyperlink, then individual task links
                links.append(("Source + Video: ", "", ""))
                
                # Create separate links for each folder so each gets its own hyperlink
                all_tasks = task.get('_all_tasks', [])
                
                for individual_task in all_tasks:
                    # Extract folder name
                    folder = individual_task.get('folder', '')
                    if isinstance(folder, str):
                        folder_name = Path(folder).name
                    else:
                        folder_name = folder.name if hasattr(folder, 'name') else str(folder)
                    
                    # Remove date prefix (e.g., "1017 ") from folder name
                    folder_name = re.sub(r'^\d{4}\s+', '', folder_name)
                    
                    # Get source link for this specific task
                    task_source_link = individual_task.get('source_video_link', '')
                    
                    # Add individual link (no prefix, just the folder name)
                    links.append(("", folder_name, task_source_link or ""))
        else:
            # Single task or base-folder API - use single source video link
            source_link = self.config.get('source_video_link', '') if self.config.get('source_video_link', '') else task.get('source_video_link', '')
            links.append(("Source + Video: ", "Link", source_link))

        for i, (pre, txt, url) in enumerate(links):
            para = info_box.text_frame.paragraphs[0] if i == 0 else info_box.text_frame.add_paragraph()
            if url:
                para.clear()
                r1, r2 = para.add_run(), para.add_run()
                r1.text, r1.font.size = pre, Pt(24)
                r2.text, r2.font.size = txt, Pt(24)
                r2.hyperlink.address = url
                para.alignment = PP_ALIGN.CENTER
            else:
                para.text, para.font.size, para.alignment = f"{pre}{txt}", Pt(24), PP_ALIGN.CENTER

    
    def _remove_template_slides(self, ppt):
        """Remove dummy template slides (slides 3 and 4) used for generation.
        
        These slides are placeholder slides in the template that are used as
        layout references for creating content slides, but should not appear
        in the final presentation.
        """
        # Slides are 0-indexed internally, so slides 3 and 4 are indices 2 and 3
        # We need to remove them in reverse order to maintain correct indices
        slides_to_remove = [3, 2]  # Remove slide 4 first, then slide 3
        
        for slide_idx in slides_to_remove:
            if len(ppt.slides) > slide_idx:
                slide_id = ppt.slides._sldIdLst[slide_idx].rId
                ppt.part.drop_rel(slide_id)
                del ppt.slides._sldIdLst[slide_idx]
    
    def save_presentation(self, ppt, task, use_comparison, effect_names=None):
        """Save the presentation"""
        try:
            # Remove dummy template slides before saving
            self._remove_template_slides(ppt)
            
            # Generate filename
            if self.api_name in ["vidu_effects", "vidu_reference", "pixverse", "kling_effects"]:
                folder_name = Path(self.config.get('base_folder', '')).name
            elif self.api_name == "veo_itv":
                # For veo_itv, use parent folder (contains date like "0130 6 Styles")
                if task.get('_is_grouped'):
                    folder_name = task
                else:
                    folder_path = Path(task.get('folder', ''))
                    folder_name = folder_path.parent.name
            else:
                # Handle grouped tasks
                if task.get('_is_grouped'):
                    folder_name = task
                else:
                    folder_name = task.get('folder', Path(self.config.get('base_folder', '')).name)
                    if isinstance(folder_name, str):
                        folder_name = Path(folder_name).name
            
            api_display = self._api_display_names.get(self.api_name, self.api_name.title())
            
            if use_comparison and task.get('reference_folder'):
                ref_name = Path(task['reference_folder']).name
                api_line, styles_line = self.get_cmp_filename(folder_name, ref_name, api_display, effect_names1=effect_names)
            else:
                api_line, styles_line = self.get_filename(folder_name, api_display, effect_names=effect_names)
            
            # Combine for filename (join the two lines with a space)
            filename = f"{api_line} {styles_line}"
            
            # Get output directory
            output_dir = Path(self.config.get('output_directory',
                            self.config.get('output', {}).get('directory',
                            self.report_definitions.get('output_directory', './'))))
            # Ensure output directory exists
            output_dir.mkdir(parents=True, exist_ok=True)
            output_path = output_dir / f"{filename}.pptx"
            # Save
            ppt.save(str(output_path))
            logger.info(f"✓ Saved: {output_path}")
            return True
        except Exception as e:
            logger.error(f"✗ Save failed: {e}")
            return False
    
    # ================== MAIN EXECUTION ==================
    
    def run(self) -> bool:
        """Main execution using unified system"""
        logger.info(f"🎬 Starting {self.api_name.title()} Report Generator")
        try:
            # Check for task grouping configuration (check both locations for backward compatibility)
            group_tasks_by = self.config.get('output', {}).get('group_tasks_by', 0) or self.config.get('group_tasks_by', 0)
            
            # Get tasks
            tasks = self.config.get('tasks', [])
            
            # Determine processing mode
            if self.api_name in ["vidu_effects", "vidu_reference", "pixverse", "kling_effects", "veo_itv"]:
                # Base folder structure APIs
                if group_tasks_by and group_tasks_by > 1 and tasks:
                    # Base-folder APIs with grouping - process tasks individually
                    logger.info(f"📦 Enabling grouped mode for {self.api_name}")
                    return self._run_grouped(tasks, group_tasks_by)
                else:
                    # Original behavior: single report for all effects
                    pairs = self.process_batch({})
                    if pairs:
                        return self.create_presentation(pairs, {})
                    else:
                        logger.warning("No pairs found for report.")
                        return False
            elif self.api_name in ["veo", "kling_ttv"]:
                # Text-to-video APIs - process once since all tasks share same root folder
                if not tasks:
                    logger.warning("No tasks found in config.")
                    return False
                
                # Process all tasks together (they all use the same root output_folder)
                # Pass the first task to get the root folder, but all videos are already there
                pairs = self.process_batch(tasks[0] if tasks else {})
                if pairs:
                    return self.create_presentation(pairs, tasks[0] if tasks else {})
                else:
                    logger.warning("No pairs found for report.")
                    return False
            else:
                # Task folder structure - multiple tasks
                if not tasks:
                    logger.warning("No tasks found in config.")
                    return False
                
                if group_tasks_by and group_tasks_by > 1:
                    # Grouped presentation mode
                    return self._run_grouped(tasks, group_tasks_by)
                else:
                    # Original individual presentation mode
                    successful = 0
                    for i, task in enumerate(tasks, 1):
                        pairs = self.process_batch(task)
                        if pairs:
                            if self.create_presentation(pairs, task):
                                successful += 1
                    logger.info(f"✓ Generated {successful}/{len(tasks)} presentations")
                    return successful > 0
        except Exception as e:
            logger.error(f"✗ Report generation failed: {e}")
            return False
        finally:
            # Cleanup all temporary files
            self.cleanup_temp_frames()
            self.cleanup_tempfiles()  # Cleanup temporary format conversions
    
    def _run_grouped(self, tasks: List[Dict], group_size: int) -> bool:
        """Process tasks in groups, creating combined presentations
        
        Works for both folder-based APIs (nano_banana, kling, etc.) 
        and base-folder APIs (vidu_effects, etc.)
        
        Tasks are separated by use_comparison_template before grouping,
        since different templates require different presentations.
        """
        logger.info(f"📦 Grouping tasks by {group_size}")
        
        # Separate tasks by template type
        comparison_tasks = []
        regular_tasks = []
        
        slide_config = self.get_slide_config()
        force_regular = slide_config.get('override_positions', False)
        for task in tasks:
            if force_regular:
                use_comparison = False
            else:
                use_comparison = task.get('use_comparison_template', False) or bool(task.get('reference_folder'))
            if use_comparison:
                comparison_tasks.append(task)
            else:
                regular_tasks.append(task)
        
        successful_groups = 0
        total_groups = 0
        
        # Process regular tasks (non-comparison template)
        if regular_tasks:
            logger.info(f"📄 Processing {len(regular_tasks)} regular template tasks")
            regular_group_count = (len(regular_tasks) + group_size - 1) // group_size
            total_groups += regular_group_count
            
            for group_idx in range(0, len(regular_tasks), group_size):
                group_tasks = regular_tasks[group_idx:group_idx + group_size]
                group_num = (group_idx // group_size) + 1
                
                logger.info(f"\n📊 Processing regular group {group_num}/{regular_group_count} ({len(group_tasks)} tasks)")
                
                task_pairs_list = []
                group_task_info = []
                
                for task in group_tasks:
                    pairs = self.process_batch(task)
                    if pairs:
                        task_pairs_list.append({'task': task, 'pairs': pairs})
                        group_task_info.append(task)
                
                if task_pairs_list:
                    combined_task = self._create_combined_task(group_task_info, group_num, regular_group_count)
                    if self.create_grouped_presentation(task_pairs_list, combined_task):
                        successful_groups += 1
                else:
                    logger.warning(f"⚠ Regular group {group_num} has no valid pairs")
        
        # Process comparison tasks (comparison template)
        if comparison_tasks:
            logger.info(f"📄 Processing {len(comparison_tasks)} comparison template tasks")
            comparison_group_count = (len(comparison_tasks) + group_size - 1) // group_size
            total_groups += comparison_group_count
            
            for group_idx in range(0, len(comparison_tasks), group_size):
                group_tasks = comparison_tasks[group_idx:group_idx + group_size]
                group_num = (group_idx // group_size) + 1
                
                logger.info(f"\n📊 Processing comparison group {group_num}/{comparison_group_count} ({len(group_tasks)} tasks)")
                
                task_pairs_list = []
                group_task_info = []
                
                for task in group_tasks:
                    pairs = self.process_batch(task)
                    if pairs:
                        task_pairs_list.append({'task': task, 'pairs': pairs})
                        group_task_info.append(task)
                
                if task_pairs_list:
                    combined_task = self._create_combined_task(group_task_info, group_num, comparison_group_count)
                    if self.create_grouped_presentation(task_pairs_list, combined_task):
                        successful_groups += 1
                else:
                    logger.warning(f"⚠ Comparison group {group_num} has no valid pairs")
        
        logger.info(f"\n✓ Generated {successful_groups}/{total_groups} grouped presentations")
        return successful_groups > 0
    
    def _create_combined_task(self, tasks: List[Dict], group_num: int, total_groups: int) -> Dict:
        """Create a combined task dict for grouped presentation
        
        Handles both folder-based APIs (folder key) and base-folder APIs (effect key)
        """
        if not tasks:
            return {}
        
        # Detect API type from task structure
        # Check for effect or custom_effect (for kling_effects and vidu_effects), and no folder key
        is_base_folder_api = ('effect' in tasks[0] or 'custom_effect' in tasks[0]) and 'folder' not in tasks[0]
        
        if is_base_folder_api:
            # Base folder API (vidu_effects, vidu_reference, pixverse, kling_effects)
            # Extract effect/category names for the combined title
            effect_names = []
            for task in tasks:
                # custom_effect has priority over effect
                effect_name = task.get('custom_effect') or task.get('effect', task.get('category', 'Unknown'))
                effect_names.append(effect_name)
            
            # Use config base folder and add grouped information
            combined = {
                'base_folder': self.config.get('base_folder', ''),
                'design_link': self.config.get('design_link', ''),
                'source_video_link': self.config.get('source_video_link', ''),
                '_is_grouped': True,
                '_is_base_folder_api': True,
                '_group_number': group_num,
                '_total_groups': total_groups,
                '_effect_names': effect_names,
                '_all_tasks': tasks
            }
        else:
            # Folder-based API (nano_banana, kling, runway, genvideo, veo_itv)
            # Extract folder names for the combined title
            folder_names = []
            parent_folder_name = None  # For veo_itv, store parent folder for date extraction
            for task in tasks:
                folder = task.get('folder', '')
                if isinstance(folder, str):
                    folder_path = Path(folder)
                    folder_name = folder_path.name
                    # For veo_itv, capture parent folder name (contains date like "0130 6 Styles")
                    if self.api_name == "veo_itv" and parent_folder_name is None:
                        parent_folder_name = folder_path.parent.name
                else:
                    folder_name = folder.name if hasattr(folder, 'name') else str(folder)
                folder_names.append(folder_name)
            
            # Use the first task as base and add grouped information
            combined = tasks[0].copy()
            combined['_is_grouped'] = True
            combined['_is_base_folder_api'] = False
            combined['_group_number'] = group_num
            combined['_total_groups'] = total_groups
            combined['_folder_names'] = folder_names
            if parent_folder_name:
                combined['_parent_folder_name'] = parent_folder_name
            combined['_all_tasks'] = tasks
        
        return combined

def create_report_generator(api_name, config_file=None):
    """Factory function to create report generator"""
    supported_apis = ['kling', 'kling_effects', 'kling_endframe', 'kling_ttv', 'kling_motion', 'nano_banana', 'vidu_effects', 'vidu_reference', 'runway', 'genvideo', 'pixverse', 'wan', 'dreamactor', 'veo', 'veo_itv']
    if api_name not in supported_apis:
        raise ValueError(f"Unsupported API: {api_name}. Supported: {supported_apis}")
    return UnifiedReportGenerator(api_name, config_file)


def main():
    import argparse
    parser = argparse.ArgumentParser(description='Generate PowerPoint reports from API processing results')
    parser.add_argument('api_name', choices=['kling', 'kling_effects', 'kling_endframe', 'kling_ttv', 'kling_motion', 'nano_banana', 'vidu_effects', 'vidu_reference', 'runway', 'genvideo', 'pixverse', 'wan', 'dreamactor', 'veo', 'veo_itv'],
                       help='API type to generate report for')
    parser.add_argument('--config', '-c', help='Config file path (optional)')
    
    args = parser.parse_args()
    
    generator = create_report_generator(args.api_name, args.config)
    sys.exit(0 if generator.run() else 1)


if __name__ == "__main__":
    main()
